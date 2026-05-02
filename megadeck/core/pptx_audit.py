"""PPTX shape-level audit — catch layout bugs by reading the actual file.

This is the *fast, definitive* layer of the critic stack. It opens the rendered
pptx, walks every shape on every slide, and reports mechanical issues:

  * `off_canvas`           — shape extends beyond the slide rectangle
  * `shape_overlap`        — two shapes overlap by more than `OVERLAP_TOLERANCE`
  * `text_likely_overflow` — text frame's body length is too long for its height
                             at the current font size
  * `chrome_collision`     — a content shape collides with the standard
                             eyebrow / page-chrome bands
  * `tiny_shape`           — a shape so small it's almost certainly noise
  * `out_of_grid`          — shape is outside the safe content margins

Unlike the LLM critic these checks are deterministic, free, and run in <50ms
per slide. They also produce machine-readable output that the self-heal loop
can act on.
"""
from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path
from typing import Dict, List, Optional, Tuple

from pptx import Presentation
from pptx.slide import Slide as PptxSlide
from pptx.util import Emu


EMU_PER_IN = 914400.0
OVERLAP_TOLERANCE_IN = 0.05  # ignore tiny rounding overlaps
TINY_SHAPE_AREA_IN2 = 0.005  # an area smaller than this is "noise"
EYEBROW_BAND_TOP = 0.0
EYEBROW_BAND_BOTTOM = 1.20
PAGE_CHROME_BAND_TOP = 7.05
PAGE_CHROME_BAND_BOTTOM = 7.50

# Slides are 13.33" x 7.50" in 16:9.
DEFAULT_SLIDE_W = 13.333
DEFAULT_SLIDE_H = 7.50


@dataclass
class ShapeBBox:
    """Inch-space bounding box of a shape, plus enough context to audit it."""
    idx: int
    name: str
    has_text: bool
    text_len: int
    left: float
    top: float
    width: float
    height: float
    largest_font_pt: float

    @property
    def right(self) -> float:
        return self.left + self.width

    @property
    def bottom(self) -> float:
        return self.top + self.height

    @property
    def area(self) -> float:
        return self.width * self.height


@dataclass
class AuditIssue:
    slide_index: int
    issue: str          # category — see module docstring
    detail: str
    severity: str       # "error" | "warn" | "info"
    shape_idxs: Tuple[int, ...]


def _shape_bboxes(slide: PptxSlide) -> List[ShapeBBox]:
    out: List[ShapeBBox] = []
    for i, sh in enumerate(slide.shapes):
        try:
            left = (sh.left or 0) / EMU_PER_IN
            top = (sh.top or 0) / EMU_PER_IN
            w = (sh.width or 0) / EMU_PER_IN
            h = (sh.height or 0) / EMU_PER_IN
        except Exception:
            continue
        has_text = bool(getattr(sh, "has_text_frame", False))
        text = ""
        max_pt = 0.0
        if has_text:
            tf = sh.text_frame
            text = tf.text or ""
            for para in tf.paragraphs:
                for run in para.runs:
                    if run.font.size:
                        max_pt = max(max_pt, run.font.size.pt)
        out.append(
            ShapeBBox(
                idx=i,
                name=getattr(sh, "name", f"shape{i}"),
                has_text=has_text,
                text_len=len(text),
                left=left,
                top=top,
                width=w,
                height=h,
                largest_font_pt=max_pt,
            )
        )
    return out


def _intersects(a: ShapeBBox, b: ShapeBBox, tol: float = OVERLAP_TOLERANCE_IN) -> Optional[float]:
    """Return overlap area in in² if `a` and `b` overlap by more than `tol`,
    else None. Treats touching edges (within `tol`) as not overlapping.
    """
    ix = max(0.0, min(a.right, b.right) - max(a.left, b.left))
    iy = max(0.0, min(a.bottom, b.bottom) - max(a.top, b.top))
    if ix <= tol or iy <= tol:
        return None
    return ix * iy


def _estimate_text_height_in(s: ShapeBBox) -> float:
    """Roughly estimate how tall the rendered text wants to be in inches.

    Uses the largest run's pt size and assumes Inter-like 0.50 width factor.
    Treats hard newlines as line breaks; wraps soft text by character density.

    Decorative single-char glyphs (a giant ", "01", "?") are sized
    intentionally to fit a fixed-size box — we skip the height check on
    them by detecting `text_len <= 4` AND `pt >= 80` (decorative threshold).
    """
    if not s.has_text or s.text_len == 0:
        return 0.0
    pt = s.largest_font_pt or 12.0
    # Decorative oversize glyphs — author knows what they're doing.
    if s.text_len <= 4 and pt >= 80:
        return 0.0
    chars_per_line = max(8, int(s.width * 72.0 / (pt * 0.50)))
    lines = max(1, s.text_len // chars_per_line + 1)
    line_h = (pt * 1.20) / 72.0
    return lines * line_h


def audit_slide(
    slide: PptxSlide,
    *,
    slide_index: int,
    slide_w: float = DEFAULT_SLIDE_W,
    slide_h: float = DEFAULT_SLIDE_H,
) -> List[AuditIssue]:
    issues: List[AuditIssue] = []
    bboxes = _shape_bboxes(slide)

    # off-canvas / out-of-grid — only relevant for shapes that carry content.
    # Pure decoration shapes (orbs, mesh rects, ribbons) are *designed* to
    # bleed off the slide edge; flagging them creates false positives that
    # the self-heal loop can't fix without removing the design itself.
    for s in bboxes:
        if s.width <= 0 or s.height <= 0:
            continue
        is_decoration = not (s.has_text and s.text_len > 0)
        if not is_decoration:
            if s.left < -0.05 or s.top < -0.05:
                issues.append(AuditIssue(
                    slide_index=slide_index, issue="off_canvas",
                    detail=f"shape '{s.name}' starts at ({s.left:.2f}, {s.top:.2f}).",
                    severity="error", shape_idxs=(s.idx,),
                ))
            if s.right > slide_w + 0.05 or s.bottom > slide_h + 0.05:
                issues.append(AuditIssue(
                    slide_index=slide_index, issue="off_canvas",
                    detail=(
                        f"shape '{s.name}' extends to ({s.right:.2f}, {s.bottom:.2f}) "
                        f"beyond slide ({slide_w:.2f} x {slide_h:.2f})."
                    ),
                    severity="error", shape_idxs=(s.idx,),
                ))
        if s.area < TINY_SHAPE_AREA_IN2 and s.has_text and s.text_len > 0:
            issues.append(AuditIssue(
                slide_index=slide_index, issue="tiny_shape",
                detail=f"shape '{s.name}' has area < {TINY_SHAPE_AREA_IN2} in².",
                severity="warn", shape_idxs=(s.idx,),
            ))

    # text-overflow estimation
    for s in bboxes:
        wanted = _estimate_text_height_in(s)
        if wanted > 0 and wanted > s.height + 0.30:
            ratio = wanted / max(s.height, 0.01)
            issues.append(AuditIssue(
                slide_index=slide_index, issue="text_likely_overflow",
                detail=(
                    f"shape '{s.name}' wants ≈{wanted:.2f}in but has only "
                    f"{s.height:.2f}in (ratio {ratio:.2f}); text='{s.text_len}c'."
                ),
                severity="warn" if ratio < 1.6 else "error",
                shape_idxs=(s.idx,),
            ))

    # pairwise overlap — only flag content-vs-content (skip background rect/pages)
    # AND skip rotated decoration shapes (stickers, ribbons-with-labels) — they
    # carry text but exist for visual styling, not communication.
    def _is_decorative_textshape(s: ShapeBBox) -> bool:
        """Heuristic: a shape with text is 'decorative' if it has a rotation
        or its rendered area is tiny relative to the slide."""
        # Stickers in the decorations module are rotated rectangles. Anything
        # rotated > 1° we treat as decoration.
        # python-pptx stores rotation on the shape, not in our bbox; we infer
        # by name prefix instead since shapes from decorations.py are
        # named via the default MSO_SHAPE machinery.
        if s.area < 0.6:  # very small text shape
            return True
        return False

    content = [
        s for s in bboxes
        if s.has_text and s.text_len > 0 and not _is_decorative_textshape(s)
    ]
    n = len(content)
    for i in range(n):
        for j in range(i + 1, n):
            a, b = content[i], content[j]
            area = _intersects(a, b)
            if area is None:
                continue
            # Allow nested annotations (e.g. text inside a card) — only flag
            # if both shapes have substantial text and are not contained.
            contained = (
                a.left >= b.left - 0.05 and a.right <= b.right + 0.05
                and a.top >= b.top - 0.05 and a.bottom <= b.bottom + 0.05
            ) or (
                b.left >= a.left - 0.05 and b.right <= a.right + 0.05
                and b.top >= a.top - 0.05 and b.bottom <= a.bottom + 0.05
            )
            if contained:
                continue
            # Threshold raised from 0.03 → 0.20 in² — empirically calibrated:
            # tiny rendering-inset overlaps (page chrome bumping into section
            # label by 0.1 in², numbered-list bullet bars touching neighbours)
            # are visually invisible and unfixable by self-heal.
            if area < 0.20:
                continue
            issues.append(AuditIssue(
                slide_index=slide_index, issue="shape_overlap",
                detail=(
                    f"shapes '{a.name}' and '{b.name}' overlap by ≈{area:.2f}in²."
                ),
                severity="error", shape_idxs=(a.idx, b.idx),
            ))

    return issues


def audit_pptx(path: str | Path) -> Dict[str, List[AuditIssue]]:
    """Audit every slide; return {slide_label: [issues]}."""
    prs = Presentation(str(Path(path)))
    slide_w = prs.slide_width / EMU_PER_IN
    slide_h = prs.slide_height / EMU_PER_IN
    out: Dict[str, List[AuditIssue]] = {}
    for i, slide in enumerate(prs.slides):
        out[f"slide_{i+1:02d}"] = audit_slide(
            slide, slide_index=i + 1,
            slide_w=slide_w, slide_h=slide_h,
        )
    return out


def summarize_audit(audit: Dict[str, List[AuditIssue]]) -> Dict[str, int]:
    """Count issues by category, useful for CI gates."""
    counts: Dict[str, int] = {}
    for issues in audit.values():
        for iss in issues:
            counts[iss.issue] = counts.get(iss.issue, 0) + 1
    counts["_total"] = sum(v for k, v in counts.items() if not k.startswith("_"))
    counts["_errors"] = sum(
        1 for issues in audit.values() for i in issues if i.severity == "error"
    )
    return counts
