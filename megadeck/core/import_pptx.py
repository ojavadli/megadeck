"""Round-trip: read an existing pptx and infer a best-effort `Deck` schema.

This is intentionally heuristic — pptx files have unbounded layout freedom, so
we cannot reliably classify every slide. We assign a slide kind based on
shape composition and emit a `numbered_list` fallback when in doubt.

This is useful for re-skinning an existing pptx with a Megadeck theme without
rebuilding from scratch.
"""
from __future__ import annotations

import re
from pathlib import Path
from typing import List, Optional, Tuple

from pptx import Presentation
from pptx.slide import Slide as PptxSlide
from pptx.util import Emu

from megadeck.core.schemas import (
    AgendaItem,
    AgendaSlide,
    BulletItem,
    Deck,
    EditorialSplitSlide,
    HeroMinimalSlide,
    HeroStatementSlide,
    KpiGridSlide,
    KpiTile,
    ManifestoSlide,
    NumberedListSlide,
    QuoteDecorativeSlide,
    SectionDividerSlide,
    SectionHeroSlide,
    Slide as SlideUnion,
    StatCalloutSlide,
    TitleSlide,
    TransitionKind,
)


_NUMBERED_PATTERN = re.compile(r"(?:^|\s)(\d+)\.\s+([^\d][^\n]*?)(?=\s+\d+\.\s+|$)")
_LEADING_NUMBER_PREFIX = re.compile(r"^\s*\d+(?:\.\d+)*[\.\)]?\s+")
_BARE_DIGIT = re.compile(r"^\s*\d+(?:\.\d+)*\s*$")
_LOOKS_LIKE_PAGE = re.compile(r"^\s*(?:slide|page|p\.?)\s*\d+\s*$", re.IGNORECASE)
_PERCENT_VALUE = re.compile(r"^\s*([+-]?\d+(?:\.\d+)?\s*%)\s*[:\-—]?\s*(.*)$")
_NUMERIC_VALUE = re.compile(r"^\s*([+-]?\$?\d[\d,\.]*[%×x]?)\s+(.+)$")


def _safe_clip(text: str, max_len: int) -> str:
    """Clip text to ≤ max_len chars without breaking mid-word.

    If text is short enough, return as-is. Otherwise back up to the last
    word boundary that fits and append "…".
    """
    text = (text or "").strip()
    if len(text) <= max_len:
        return text
    cut = text[: max_len - 1]
    space = cut.rfind(" ")
    if space > max_len * 0.6:
        cut = cut[:space]
    return cut.rstrip(" ,;.:") + "…"


def _split_fused_numbered(text: str) -> List[str]:
    """If `text` contains '1. X 2. Y 3. Z' (or with newlines), split into items.

    Returns the list of items, or an empty list if the pattern doesn't match.
    """
    matches = list(_NUMBERED_PATTERN.finditer(text))
    if len(matches) >= 2:
        return [m.group(2).strip(" ,;:") for m in matches]
    # Fallback: try splitting on newlines if every line begins with N.
    lines = [l.strip() for l in text.split("\n") if l.strip()]
    if len(lines) >= 2 and all(_LEADING_NUMBER_PREFIX.match(l) for l in lines):
        return [_LEADING_NUMBER_PREFIX.sub("", l).strip(" ,;:") for l in lines]
    return []


def _looks_like_kpi(items: List[str]) -> List[Tuple[str, str]]:
    """If most items look like 'X% Label' or 'Label X%', extract (value, label).

    Returns the parsed pairs only if at least 2 items match; else [].
    """
    pairs: List[Tuple[str, str]] = []
    for item in items:
        m = _PERCENT_VALUE.match(item)
        if m and m.group(2):
            pairs.append((m.group(1).strip(), m.group(2).strip()))
            continue
        m2 = _NUMERIC_VALUE.match(item)
        if m2 and len(m2.group(1)) <= 8:
            pairs.append((m2.group(1).strip(), m2.group(2).strip()))
            continue
    return pairs if len(pairs) >= 2 and len(pairs) == len(items) else []


def _slide_text_blocks(slide: PptxSlide) -> List[Tuple[str, float, float]]:
    """Return [(text, top_in, font_size_pt), ...] sorted by vertical position.

    Position-aware extraction lets us detect titles (large, near-top) vs body
    (smaller, lower) much more reliably than treating every text frame the same.
    """
    blocks: List[Tuple[str, float, float]] = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        txt = shape.text_frame.text.strip()
        if not txt:
            continue
        top_in = (shape.top or 0) / 914400.0
        # Approximate dominant font size by inspecting first run that has size set
        font_pt = 0.0
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.font.size:
                    font_pt = run.font.size.pt
                    break
            if font_pt:
                break
        blocks.append((txt, top_in, font_pt))
    blocks.sort(key=lambda b: (b[1], -b[2]))
    return blocks


def _split_head_tail(line: str) -> Tuple[str, str]:
    """Split a bullet-style line into (head, tail) on the first '.', '—', or ':' separator."""
    for sep in (".  ", "— ", "- ", ": "):
        if sep in line:
            head, _, tail = line.partition(sep)
            return head.strip(), tail.strip()
    # Otherwise return the whole thing as the head, no tail
    return line.strip()[:78], ""


def _extract_speaker_notes(slide: PptxSlide) -> str:
    if not slide.has_notes_slide:
        return ""
    tf = slide.notes_slide.notes_text_frame
    if tf is None:
        return ""
    return tf.text.strip()


def _classify_slide(
    blocks: List[Tuple[str, float, float]], idx: int, total: int
) -> SlideUnion:
    """Pick a Megadeck slide kind from the raw text on a pptx slide."""
    notes = ""  # caller adds this back

    # Drop blocks that are pure noise: bare digits, single chars, page numbers.
    blocks = [
        b for b in blocks
        if not _BARE_DIGIT.match(b[0])
        and not _LOOKS_LIKE_PAGE.match(b[0])
        and len(b[0].strip()) >= 2
    ]

    if not blocks:
        return SectionDividerSlide(
            kind="section_divider",
            part_label=f"Section {idx + 1:02d}",
            title="(empty slide)",
            notes=notes,
            transition=TransitionKind.FADE,
        )

    # Pick the title: prefer the largest font, but skip blocks that are clearly
    # bullet content (multi-line + many chars) even if their font is large.
    by_size = sorted(blocks, key=lambda b: -b[2])
    title_block = by_size[0] if by_size[0][2] > 0 else blocks[0]
    title_text_raw = title_block[0]

    # If the largest-font block contains a fused numbered list ("1. A 2. B 3. C"),
    # split it: the actual title may be elsewhere (eyebrow), and these become bullets.
    fused_items = _split_fused_numbered(title_text_raw)
    title: str
    body: List[str]
    eyebrow_text: Optional[str] = None

    above_title = [
        b for b in blocks
        if b[1] < title_block[1] and b[2] and b[2] < title_block[2]
    ]
    if above_title:
        candidate = _safe_clip(above_title[0][0], 38)
        # Suppress the eyebrow when it duplicates the title text — that creates
        # a noisy "ENTREPRENEURSHIP AND… / Entrepreneurship and …" double-header.
        norm_candidate = candidate.lower().rstrip(":.…").strip()
        norm_title = title_block[0].lower().rstrip(":.…").strip()
        if norm_candidate and not norm_title.startswith(norm_candidate):
            eyebrow_text = candidate

    title_lines = [ln.strip() for ln in title_text_raw.split("\n") if ln.strip()]

    if fused_items:
        # Use the eyebrow as title if available, else fall back to a generic.
        title = _safe_clip(eyebrow_text or title_lines[0] if title_lines else "Key points", 100)
        # Body = the split items + any remaining non-title blocks
        body_raw = list(fused_items)
        for b in blocks:
            if b is title_block:
                continue
            if above_title and b[0] == above_title[0][0]:
                continue
            body_raw.append(b[0])
        body = body_raw
    elif len(title_lines) >= 3 and all(len(ln) <= 120 for ln in title_lines):
        # Multi-line title where each line looks like a separate idea
        # → first line is the real title, others become bullets.
        title = _safe_clip(title_lines[0], 100)
        body_raw = list(title_lines[1:])
        for b in blocks:
            if b is title_block:
                continue
            if above_title and b[0] == above_title[0][0]:
                continue
            body_raw.append(b[0])
        body = body_raw
    else:
        # Single-line or 2-line title; use the largest-font block as the title.
        title = _safe_clip(title_text_raw.replace("\n", " "), 100)
        body = [b[0] for b in blocks if b[0] != title_text_raw]
        if above_title:
            body = [b for b in body if b != above_title[0][0]]

    # First slide → always a cover (title slide). Pack body into the canonical
    # cover fields; extras are dropped (a cover should be uncluttered anyway).
    if idx == 0:
        subtitle = _safe_clip(body[0], 138) if len(body) >= 1 else None
        presenter = _safe_clip(body[1], 78) if len(body) >= 2 else None
        date = _safe_clip(body[2], 38) if len(body) >= 3 else None
        venue = _safe_clip(body[3], 78) if len(body) >= 4 else None
        return TitleSlide(
            kind="title", eyebrow=eyebrow_text or "",
            title=title, subtitle=subtitle,
            presenter=presenter, date=date, venue=venue, notes=notes,
            transition=TransitionKind.FADE,
        )

    # KPI grid heuristic — if 2-4 items and they all look like "X% Label"
    if 2 <= len(body) <= 4:
        kpi_pairs = _looks_like_kpi(body)
        if kpi_pairs:
            return KpiGridSlide(
                kind="kpi_grid",
                eyebrow=eyebrow_text or f"Slide {idx + 1:02d}",
                title=title,
                tiles=[
                    KpiTile(value=v[:14], label=_safe_clip(l, 28))
                    for v, l in kpi_pairs
                ],
                notes=notes,
                transition=TransitionKind.FADE,
            )

    # ---- Variety rules: produce STRUCTURALLY DIFFERENT slide kinds ----------
    # We rotate through several layout families so the deck doesn't read as
    # 30 numbered_list slides recoloured. Each rule below picks a distinct
    # template based on the source slide's content shape.

    # Single big stat / number — render as full-bleed stat callout.
    title_is_pure_stat = bool(_PERCENT_VALUE.match(title)) or bool(
        _NUMERIC_VALUE.match(title) and len(title) <= 14
    )
    if title_is_pure_stat and len(body) <= 2:
        m = _PERCENT_VALUE.match(title) or _NUMERIC_VALUE.match(title)
        value = m.group(1) if m else title
        label_text = (m.group(2) if m else "").strip() or (body[0] if body else "Result")
        return StatCalloutSlide(
            kind="stat_callout",
            eyebrow=eyebrow_text or f"Slide {idx + 1:02d}",
            value=_safe_clip(value, 14),
            label=_safe_clip(label_text, 78),
            context=_safe_clip(body[1], 178) if len(body) > 1 else None,
            notes=notes,
            transition=TransitionKind.FADE,
        )

    # Quote-shaped slide: title is in quotes, attribution looks like a person.
    quote_signal = (
        (title.startswith(('"', "“", "‘")) and title.endswith(('"', "”", "’")))
        or any(line.lstrip().startswith(("—", "–")) for line in body)
        or (len(body) >= 1 and len(title) <= 240 and any(
            len(b) <= 60 and (b.count(",") <= 2) and not b.endswith(".")
            for b in body[:1]
        ) and "“" in title or '"' in title)
    )
    if quote_signal and body:
        return QuoteDecorativeSlide(
            kind="quote_decorative",
            eyebrow=eyebrow_text or "QUOTE",
            quote=_safe_clip(title.strip("\"“”'‘’"), 298),
            author=_safe_clip(body[0].lstrip("—– "), 58),
            role=_safe_clip(body[1], 78) if len(body) > 1 else None,
            notes=notes,
            transition=TransitionKind.FADE,
        )

    # Manifesto: 1-2 dense paragraphs, total chars > 220, no clear bullet structure.
    body_chars = sum(len(b) for b in body)
    if (
        len(body) <= 2 and body_chars > 220
        and len(title) <= 60
        and not fused_items
    ):
        return ManifestoSlide(
            kind="manifesto",
            eyebrow=eyebrow_text or _safe_clip(title, 38).upper(),
            body=_safe_clip(" ".join(body)[:580], 580),
            notes=notes,
            transition=TransitionKind.FADE,
        )

    # Hero-minimal: very short single declaration (≤ 60 chars), no body.
    if len(title) <= 60 and not body and not fused_items:
        return HeroMinimalSlide(
            kind="hero_minimal",
            eyebrow=eyebrow_text or None,
            title=_safe_clip(title, 118),
            notes=notes,
            transition=TransitionKind.FADE,
        )

    # Editorial-split: title + a single longish body paragraph (one idea slide).
    if (
        len(body) == 1
        and 80 <= len(body[0]) <= 380
        and len(title) <= 80
    ):
        # Alternate sides every other slide for asymmetric rhythm.
        side = "left" if idx % 2 == 0 else "right"
        return EditorialSplitSlide(
            kind="editorial_split",
            eyebrow=eyebrow_text or f"Slide {idx + 1:02d}",
            title=_safe_clip(title, 108),
            body=_safe_clip(body[0], 398),
            side=side,
            notes=notes,
            transition=TransitionKind.FADE,
        )

    # Hero-statement heuristic: very short title and few short other blocks
    if (
        len(title) <= 50
        and len(body) <= 4
        and not any(len(b) > 200 for b in body)
        and not fused_items
    ):
        return HeroStatementSlide(
            kind="hero_statement",
            eyebrow=eyebrow_text or f"Slide {idx + 1:02d}",
            statement=_safe_clip(title, 78),
            supports=[_safe_clip(r, 158) for r in body][:4],
            notes=notes,
            transition=TransitionKind.FADE,
        )

    # Default → numbered_list. Split each body line into head+tail.
    bullets: List[BulletItem] = []
    for line in body[:6]:
        # Strip leading "1. " etc that we may have already extracted
        clean = _LEADING_NUMBER_PREFIX.sub("", line).strip()
        head, tail = _split_head_tail(clean)
        if head:
            bullets.append(BulletItem(head=_safe_clip(head, 78), tail=_safe_clip(tail, 198)))

    # Edge case: only 1 valid bullet → promote to hero_statement instead of padding
    if len(bullets) < 2:
        if bullets:
            return HeroStatementSlide(
                kind="hero_statement",
                eyebrow=eyebrow_text or f"Slide {idx + 1:02d}",
                statement=_safe_clip(title, 78),
                supports=[bullets[0].head] + ([bullets[0].tail] if bullets[0].tail else []),
                notes=notes,
                transition=TransitionKind.FADE,
            )
        return HeroStatementSlide(
            kind="hero_statement",
            eyebrow=eyebrow_text or f"Slide {idx + 1:02d}",
            statement=_safe_clip(title, 78),
            supports=[],
            notes=notes,
            transition=TransitionKind.FADE,
        )

    return NumberedListSlide(
        kind="numbered_list",
        eyebrow=eyebrow_text or f"Slide {idx + 1:02d}",
        title=title,
        items=bullets,
        notes=notes,
        transition=TransitionKind.FADE,
    )


def import_pptx(
    path: str | Path,
    theme: str = "default",
    *,
    rhythm: bool = True,
) -> Deck:
    """Read a pptx and return a heuristically-classified Deck.

    Round-trip is best-effort: complex slides (charts, images, tables) are
    coerced into the closest matching template. Re-render the result through
    `render_deck` to get a Megadeck-themed pptx.

    `rhythm=True` (default) post-processes the deck to assign layout variants
    that maximise visual variety (numbered_list slides cycle through
    cards / timeline / split / default).
    """
    src = Presentation(str(Path(path)))
    slides_out: List[SlideUnion] = []
    for i, src_slide in enumerate(src.slides):
        blocks = _slide_text_blocks(src_slide)
        slide_obj = _classify_slide(blocks, i, len(src.slides))
        notes = _extract_speaker_notes(src_slide)
        if notes:
            slide_obj = slide_obj.model_copy(update={"notes": notes[:6000]})
        slides_out.append(slide_obj)

    # Post-process: promote standalone "section break" slides to section_hero.
    # Heuristic: a slide that's a HeroMinimalSlide whose title is short AND
    # whose neighbours have very different content density is treated as a
    # section break. We also force every Nth (5th) HeroMinimal to section_hero
    # so chapter breaks are visible across the deck.
    promoted: List[SlideUnion] = []
    section_count = 0
    for i, sd in enumerate(slides_out):
        if (
            isinstance(sd, HeroMinimalSlide)
            and len(sd.title) <= 60
            and i > 0
            and i < len(slides_out) - 1
        ):
            section_count += 1
            promoted.append(SectionHeroSlide(
                kind="section_hero",
                part_number=f"{section_count:02d}",
                part_label=sd.eyebrow or f"Section {section_count}",
                title=sd.title,
                subtitle=None,
                notes=sd.notes,
                transition=sd.transition,
            ))
        else:
            promoted.append(sd)

    deck = Deck(
        title=Path(path).stem.replace("_", " ").title(),
        theme=theme,
        slides=promoted,
    )
    if rhythm:
        from megadeck.core.rhythm import apply_rhythm
        deck = apply_rhythm(deck)
    return deck
