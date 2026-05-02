"""Renderer — turns a validated `Deck` schema into a `.pptx` file."""
from __future__ import annotations

from pathlib import Path
from typing import Optional

from lxml import etree
from pptx import Presentation
from pptx.util import Emu, Inches

from megadeck.animations.transitions import apply_transition
from megadeck.core.schemas import (
    AgendaSlide,
    BeforeAfterSlide,
    BentoGridSlide,
    CallToActionSlide,
    CodeSnippetSlide,
    ComparisonTableSlide,
    Deck,
    EditorialSplitSlide,
    FaqListSlide,
    FeatureGridSlide,
    HeroMinimalSlide,
    HeroStatementSlide,
    KpiGridSlide,
    LogoGridSlide,
    ManifestoSlide,
    NumberedListSlide,
    PricingTableSlide,
    PullQuoteSlide,
    QuestionSlide,
    QuoteDecorativeSlide,
    SectionDividerSlide,
    SectionHeroSlide,
    StatCalloutSlide,
    StatHeroSlide,
    StepDiagramSlide,
    SwotMatrixSlide,
    TakeawaysSlide,
    TeamGridSlide,
    TestimonialGridSlide,
    ThreeCardSlide,
    TimelineSlide,
    TitleSlide,
    TwoColumnSlide,
)
from megadeck.design_system.templates import (
    render_agenda,
    render_before_after,
    render_bento_grid,
    render_call_to_action,
    render_code_snippet,
    render_comparison_table,
    render_editorial_split,
    render_faq_list,
    render_feature_grid,
    render_hero_minimal,
    render_hero_statement,
    render_kpi_grid,
    render_logo_grid,
    render_manifesto,
    render_numbered_list,
    render_pricing_table,
    render_pull_quote,
    render_question,
    render_quote_decorative,
    render_section_divider,
    render_section_hero,
    render_stat_callout,
    render_stat_hero,
    render_step_diagram,
    render_swot_matrix,
    render_takeaways,
    render_team_grid,
    render_testimonial_grid,
    render_three_card,
    render_timeline,
    render_title,
    render_two_column,
)
from megadeck.design_system.tokens import Theme, get_theme


_RENDERERS = {
    HeroStatementSlide: render_hero_statement,
    NumberedListSlide: render_numbered_list,
    ThreeCardSlide: render_three_card,
    TwoColumnSlide: render_two_column,
    SectionDividerSlide: render_section_divider,
    AgendaSlide: render_agenda,
    TimelineSlide: render_timeline,
    ComparisonTableSlide: render_comparison_table,
    PullQuoteSlide: render_pull_quote,
    BentoGridSlide: render_bento_grid,
    KpiGridSlide: render_kpi_grid,
    BeforeAfterSlide: render_before_after,
    StepDiagramSlide: render_step_diagram,
    CodeSnippetSlide: render_code_snippet,
    TitleSlide: render_title,
    FeatureGridSlide: render_feature_grid,
    TestimonialGridSlide: render_testimonial_grid,
    TeamGridSlide: render_team_grid,
    PricingTableSlide: render_pricing_table,
    TakeawaysSlide: render_takeaways,
    CallToActionSlide: render_call_to_action,
    SwotMatrixSlide: render_swot_matrix,
    FaqListSlide: render_faq_list,
    StatHeroSlide: render_stat_hero,
    LogoGridSlide: render_logo_grid,
    QuestionSlide: render_question,
    HeroMinimalSlide: render_hero_minimal,
    EditorialSplitSlide: render_editorial_split,
    StatCalloutSlide: render_stat_callout,
    ManifestoSlide: render_manifesto,
    QuoteDecorativeSlide: render_quote_decorative,
    SectionHeroSlide: render_section_hero,
}


P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
NSMAP = {"a": A_NS, "p": P_NS}


def _set_speaker_notes(slide, text: str) -> None:
    """Inject speaker notes into a slide; create the placeholder if needed."""
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    if tf is not None:
        tf.text = text
        return
    # Fallback path — fabricate the body placeholder
    notes_root = notes_slide.element
    body_sp = None
    for sp in notes_root.findall(".//p:sp", NSMAP):
        ph = sp.find(".//p:ph", NSMAP)
        if ph is not None and (ph.get("type") == "body" or ph.get("idx") == "3"):
            body_sp = sp
            break
    if body_sp is None:
        sp_xml = (
            f'<p:sp xmlns:a="{A_NS}" xmlns:p="{P_NS}">'
            '<p:nvSpPr><p:cNvPr id="3" name="Notes Placeholder 2"/>'
            '<p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr>'
            '<p:nvPr><p:ph type="body" sz="quarter" idx="3"/></p:nvPr>'
            "</p:nvSpPr><p:spPr/></p:sp>"
        )
        body_sp = etree.fromstring(sp_xml)
        notes_root.find(".//p:spTree", NSMAP).append(body_sp)
    existing = body_sp.find("p:txBody", NSMAP)
    if existing is not None:
        body_sp.remove(existing)
    tx_body = etree.SubElement(body_sp, "{%s}txBody" % P_NS)
    etree.SubElement(tx_body, "{%s}bodyPr" % A_NS)
    etree.SubElement(tx_body, "{%s}lstStyle" % A_NS)
    for line in text.split("\n"):
        p = etree.SubElement(tx_body, "{%s}p" % A_NS)
        if line:
            r = etree.SubElement(p, "{%s}r" % A_NS)
            r_pr = etree.SubElement(r, "{%s}rPr" % A_NS)
            r_pr.set("lang", "en-US")
            r_pr.set("sz", "1100")
            t = etree.SubElement(r, "{%s}t" % A_NS)
            t.text = line


def _section_label_for(deck: Deck, idx: int) -> Optional[str]:
    """Return the most recent section divider's label (used in the page chrome)."""
    label: Optional[str] = None
    for j, slide in enumerate(deck.slides):
        if j > idx:
            break
        if isinstance(slide, SectionDividerSlide):
            label = f"{slide.part_label.upper()}  ·  {slide.title.upper()}"
    return label


def render_deck(deck: Deck, output_path: str | Path) -> Path:
    """Render a validated `Deck` to a `.pptx` file. Returns the output path."""
    out = Path(output_path)
    out.parent.mkdir(parents=True, exist_ok=True)

    theme: Theme = get_theme(deck.theme)

    prs = Presentation()
    # 16:9 widescreen
    prs.slide_width = Emu(int(theme.slide_width_in * 914400))
    prs.slide_height = Emu(int(theme.slide_height_in * 914400))

    # Use the blank layout for everything — we draw all chrome ourselves.
    blank_layout = prs.slide_layouts[6]

    total = len(deck.slides)
    for idx, sdata in enumerate(deck.slides):
        slide = prs.slides.add_slide(blank_layout)
        page_n = idx + 1
        section_label = _section_label_for(deck, idx)

        renderer_fn = _RENDERERS.get(type(sdata))
        if renderer_fn is None:
            raise NotImplementedError(f"Unknown slide kind: {type(sdata).__name__}")
        # Variant dispatch: a slide may opt into a layout variant of its kind
        # (e.g. numbered_list 'split' instead of the default outlined-numerals).
        # Themes can default-lock variants too. See megadeck.design_system.variants.
        from megadeck.design_system.variants import get_variant_renderer
        renderer_fn = get_variant_renderer(
            kind=getattr(sdata, "kind", ""),
            variant=getattr(sdata, "variant", None),
            theme=theme,
            default=renderer_fn,
        )
        renderer_fn(
            slide, sdata, theme,
            page_n=page_n, page_total=total,
            section_label=section_label,
        )

        # Speaker notes
        if sdata.notes:
            _set_speaker_notes(slide, sdata.notes)

        # Slide transition
        apply_transition(slide.element, sdata.transition)

    prs.save(str(out))
    return out
