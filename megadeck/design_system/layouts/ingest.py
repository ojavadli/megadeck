"""Ingest .pptx files into Megadeck layouts.

For each slide we walk every shape and produce a `LayoutShape` describing
geometry + role + style. Role assignment uses a small heuristic:

* role inferred from font size (largest = title; mid = subtitle/body)
* very small text near edges → eyebrow / page-chrome → 'unknown'
* picture shapes → 'image' / 'icon' depending on aspect ratio
* tall thin rect with no text → 'accent_bar'
* round small numbered text → 'number'
* large filled rect with no text → 'decoration'

The output is saved as a JSON layout file that the fill renderer can
re-instantiate with new content.
"""
from __future__ import annotations

import json
import re
from pathlib import Path
from typing import List, Optional, Tuple

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.shapes.base import BaseShape
from pptx.util import Emu

from megadeck.design_system.layouts.registry import (
    Layout,
    LayoutShape,
    register_layout_obj,
    save_layout_json,
)


_EMU_PER_INCH = 914400


def _emu_to_in(emu: int) -> float:
    return emu / _EMU_PER_INCH


def _shape_text(shape: BaseShape) -> str:
    if not shape.has_text_frame:
        return ""
    parts: List[str] = []
    for p in shape.text_frame.paragraphs:
        for r in p.runs:
            if r.text:
                parts.append(r.text)
    return "".join(parts).strip()


def _largest_run_font_size(shape: BaseShape) -> Optional[float]:
    if not shape.has_text_frame:
        return None
    best: Optional[float] = None
    for p in shape.text_frame.paragraphs:
        for r in p.runs:
            if r.font.size is not None:
                size_pt = r.font.size.pt
                if best is None or size_pt > best:
                    best = size_pt
    return best


def _font_attrs(shape: BaseShape) -> Tuple[Optional[str], bool, bool, Optional[str]]:
    """Pick the dominant font from the shape — name, bold, italic, color."""
    if not shape.has_text_frame:
        return None, False, False, None
    font_name: Optional[str] = None
    bold = False
    italic = False
    color: Optional[str] = None
    for p in shape.text_frame.paragraphs:
        for r in p.runs:
            if r.font.name and not font_name:
                font_name = r.font.name
            if r.font.bold:
                bold = True
            if r.font.italic:
                italic = True
            try:
                if r.font.color and r.font.color.rgb:
                    color = str(r.font.color.rgb)
                    break
            except Exception:
                pass
        if color:
            break
    return font_name, bold, italic, color


def _shape_fill_hex(shape: BaseShape) -> Optional[str]:
    try:
        fill = shape.fill
        if fill.type == 1:  # MSO_FILL.SOLID
            return str(fill.fore_color.rgb)
    except Exception:
        return None
    return None


def _classify_shape_role(
    shape: BaseShape,
    *,
    text: str,
    font_size_pt: Optional[float],
    width_in: float,
    height_in: float,
    top_in: float,
    is_picture: bool,
) -> str:
    """Heuristic role assignment."""
    if is_picture:
        # Square-ish small images = icon, wider/taller = photo.
        if max(width_in, height_in) <= 1.5 and abs(width_in - height_in) < 0.6:
            return "icon"
        return "image"
    if not text:
        # Thin tall rect = accent bar
        if width_in <= 0.25 and height_in >= 0.6:
            return "accent_bar"
        # Rest = decoration
        return "decoration"
    # Text-bearing roles by font size
    if font_size_pt is None:
        font_size_pt = 14.0
    txt_strip = text.strip()
    txt_short = len(txt_strip) <= 18
    # A small numeric leading character is probably a numbered-list marker
    if (
        txt_short
        and re.match(r"^\s*(\d{1,2}|0?\d|[IVX]{1,5})\s*\.?\s*$", txt_strip)
    ):
        return "number"
    # Tiny text near top edges = eyebrow / kicker
    if font_size_pt <= 14 and top_in <= 1.5 and txt_short:
        return "eyebrow"
    if font_size_pt >= 36:
        return "title"
    if font_size_pt >= 22:
        return "subtitle"
    return "body"


def _alignment_from_para(shape: BaseShape) -> str:
    if not shape.has_text_frame:
        return "left"
    try:
        for p in shape.text_frame.paragraphs:
            if p.alignment is not None:
                a = str(p.alignment)
                if "CENTER" in a:
                    return "center"
                if "RIGHT" in a:
                    return "right"
                return "left"
    except Exception:
        pass
    return "left"


def _layout_kind_hint(shapes: List[LayoutShape]) -> str:
    """Best-guess slide kind from shape composition."""
    titles = [s for s in shapes if s.role == "title"]
    bodies = [s for s in shapes if s.role == "body"]
    numbers = [s for s in shapes if s.role == "number"]
    images = [s for s in shapes if s.role == "image"]
    n_title = len(titles)
    n_body = len(bodies)
    n_num = len(numbers)
    n_img = len(images)

    if n_num >= 3:
        return "numbered_list"
    if n_img >= 1 and n_title <= 1 and n_body <= 1:
        return "photo_card"
    if n_title == 1 and n_body == 0 and n_num == 0:
        # Hero — single big title
        title = titles[0]
        if title.font_size_pt and title.font_size_pt >= 60:
            return "hero_statement"
        return "title"
    if n_body >= 4 or sum(len(s.text or "") for s in bodies) > 600:
        return "manifesto"
    if n_body >= 1:
        return "two_column"
    return "unknown"


def ingest_pptx(
    path: str | Path,
    *,
    save_to: Optional[Path] = None,
    name_prefix: Optional[str] = None,
) -> List[Layout]:
    """Walk a .pptx and emit one Layout per slide.

    Layouts are registered in memory and (optionally) saved as JSON to
    `save_to / <prefix>_<idx>.json`.
    """
    p = Path(path)
    pres = Presentation(str(p))
    slide_w_in = _emu_to_in(pres.slide_width)
    slide_h_in = _emu_to_in(pres.slide_height)

    out: List[Layout] = []
    src_name = name_prefix or p.stem.replace(" ", "_")[:40]

    for idx, src_slide in enumerate(pres.slides):
        shapes: List[LayoutShape] = []
        for z, sh in enumerate(src_slide.shapes):
            try:
                left_in = _emu_to_in(sh.left or 0)
                top_in = _emu_to_in(sh.top or 0)
                width_in = _emu_to_in(sh.width or 0)
                height_in = _emu_to_in(sh.height or 0)
            except Exception:
                continue
            text = _shape_text(sh)
            font_size_pt = _largest_run_font_size(sh)
            font_name, bold, italic, color = _font_attrs(sh)
            fill_hex = _shape_fill_hex(sh)
            is_picture = sh.shape_type == MSO_SHAPE_TYPE.PICTURE
            role = _classify_shape_role(
                sh,
                text=text,
                font_size_pt=font_size_pt,
                width_in=width_in,
                height_in=height_in,
                top_in=top_in,
                is_picture=is_picture,
            )
            shapes.append(LayoutShape(
                role=role,
                left_in=left_in, top_in=top_in,
                width_in=width_in, height_in=height_in,
                text=text or None,
                font_size_pt=font_size_pt,
                font_bold=bool(bold),
                font_italic=bool(italic),
                font_name=font_name,
                color=color,
                fill=fill_hex,
                alignment=_alignment_from_para(sh),
                z=z,
            ))

        # Filter out micro-shapes with zero area
        shapes = [s for s in shapes if s.width_in > 0.05 and s.height_in > 0.05]
        if not shapes:
            continue

        kind_hint = _layout_kind_hint(shapes)

        # Collect a small palette of fills observed
        palette = tuple(sorted({s.fill for s in shapes if s.fill}))[:6]

        lay = Layout(
            name=f"{src_name}-{idx+1:03d}",
            source=f"file:{p.name}:{idx+1}",
            kind_hint=kind_hint,
            slide_w_in=slide_w_in,
            slide_h_in=slide_h_in,
            shapes=shapes,
            palette=palette,
            description=(
                f"Layout extracted from slide {idx+1} of {p.name}"
            ),
        )
        register_layout_obj(lay)
        out.append(lay)

        if save_to is not None:
            save_layout_json(lay, save_to / f"{lay.name}.json")

    return out


def ingest_folder(
    folder: str | Path,
    *,
    save_to: Optional[Path] = None,
    pattern: str = "*.pptx",
) -> List[Layout]:
    """Walk a folder of .pptx files and ingest every slide."""
    folder = Path(folder)
    out: List[Layout] = []
    for src in sorted(folder.rglob(pattern)):
        try:
            out.extend(ingest_pptx(src, save_to=save_to))
        except Exception as exc:  # noqa: BLE001
            import warnings
            warnings.warn(f"Failed to ingest {src}: {exc}")
    return out


__all__ = ["ingest_pptx", "ingest_folder"]
