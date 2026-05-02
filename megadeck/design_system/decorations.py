"""Decorative composition primitives — the things that make slides look *designed*.

Themes get a `decorations: List[DecorationSpec]` field. Each spec is a small
dict describing one decoration to lay down before the content renders. The
specs are validated by Pydantic and applied in order, behind the eyebrow
and content text.

Catalog (current)
-----------------

orb               Soft gradient circle, alpha-blended into bg. Defines the
                  "Aurora" / Linear / v0 look. Specify `size`, `x`, `y`,
                  `color`, `alpha`. Position is in-fraction (0-1).

mesh              Three-stop linear gradient layer over the bg, low alpha.
                  Adds the "Stripe / Vercel mesh" feel without obscuring
                  text. Specify `colors[3]`, `angle`, `alpha`.

corner_glow       Radial glow anchored to a corner ("tl"/"tr"/"bl"/"br" or
                  "top"/"bottom"). Colour bleeds 30-40% into the slide.

edge_ribbon       Thin vertical or horizontal accent stripe at a slide edge.
                  Specify `edge` ("left"/"right"/"top"/"bottom"), `color`,
                  `width_in`.

geometric         Solid-fill geometric shape (`triangle`/`diamond`/`circle`/
                  `pill`/`square`). Useful as off-corner decoration.

scribble_dots     Scatter of small dots — adds a sketchy paper feel.

aurora_band       Wide diagonal band of layered gradient orbs, used as a
                  hero accent on cinematic themes.

sticker           Filled rotated rectangle with a subtle shadow — gives a
                  "tag" / "sticker" feel for editorial themes.
"""
from __future__ import annotations

import random
from typing import List, Literal, Optional

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.slide import Slide
from pptx.util import Inches, Pt
from pydantic import BaseModel, Field

from megadeck.design_system.effects import (
    apply_drop_shadow,
    apply_linear_gradient,
    apply_radial_gradient,
    apply_solid_fill,
    apply_soft_edge,
)


# ---------------------------------------------------------------------------
# Pydantic specs — these are what themes carry in `theme.decorations`
# ---------------------------------------------------------------------------

class OrbSpec(BaseModel):
    kind: Literal["orb"] = "orb"
    x: float = Field(..., ge=-0.5, le=1.5, description="Center x as fraction of slide width.")
    y: float = Field(..., ge=-0.5, le=1.5, description="Center y as fraction of slide height.")
    size_in: float = Field(4.0, gt=0.0, le=14.0)
    color: str = Field(..., description="#RRGGBB")
    alpha_pct: int = Field(60, ge=1, le=100)
    soft_edge_pt: float = Field(20.0, ge=0.0, le=80.0)


class MeshSpec(BaseModel):
    kind: Literal["mesh"] = "mesh"
    colors: List[str] = Field(..., min_length=3, max_length=3)
    angle_deg: float = Field(135.0)
    alpha_pct: int = Field(45, ge=1, le=100)


class CornerGlowSpec(BaseModel):
    kind: Literal["corner_glow"] = "corner_glow"
    corner: Literal["tl", "tr", "bl", "br", "top", "bottom"] = "tl"
    color: str
    alpha_pct: int = Field(70, ge=1, le=100)
    radius_frac: float = Field(0.45, ge=0.1, le=0.9)


class EdgeRibbonSpec(BaseModel):
    kind: Literal["edge_ribbon"] = "edge_ribbon"
    edge: Literal["left", "right", "top", "bottom"] = "left"
    color: str
    width_in: float = Field(0.18, gt=0.0, le=0.6)
    inset_in: float = Field(0.0, ge=0.0, le=1.0)


class GeometricSpec(BaseModel):
    kind: Literal["geometric"] = "geometric"
    shape: Literal["triangle", "diamond", "circle", "pill", "square"] = "triangle"
    x: float = Field(...)
    y: float = Field(...)
    size_in: float = Field(1.5, gt=0.0)
    color: str
    alpha_pct: int = Field(100, ge=1, le=100)
    rotation_deg: float = Field(0.0)


class ScribbleDotsSpec(BaseModel):
    kind: Literal["scribble_dots"] = "scribble_dots"
    color: str
    count: int = Field(20, ge=4, le=120)
    seed: int = Field(7)
    radius_pt: float = Field(2.5, gt=0.5, le=8.0)
    alpha_pct: int = Field(30, ge=1, le=100)


class AuroraBandSpec(BaseModel):
    kind: Literal["aurora_band"] = "aurora_band"
    colors: List[str] = Field(..., min_length=2, max_length=4)
    angle_deg: float = Field(20.0)
    band_height_in: float = Field(2.5, gt=0.5, le=6.0)
    y_frac: float = Field(0.55, ge=0.0, le=1.0)


class StickerSpec(BaseModel):
    kind: Literal["sticker"] = "sticker"
    text: str = Field(..., max_length=20)
    x: float
    y: float
    width_in: float = Field(1.5)
    height_in: float = Field(0.45)
    color: str
    text_color: str = "#FFFFFF"
    rotation_deg: float = Field(-4.0)


# ---------------------------------------------------------------------------
# Renderers — each takes `slide`, the spec, and the slide dimensions in inches
# ---------------------------------------------------------------------------

def _hex(color: str) -> str:
    return color.lstrip("#").upper()


def _rgb(color: str) -> RGBColor:
    s = _hex(color)
    return RGBColor(int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16))


def _add_orb(slide: Slide, spec: OrbSpec, w: float, h: float) -> None:
    cx, cy = spec.x * w, spec.y * h
    half = spec.size_in / 2
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(cx - half), Inches(cy - half),
        Inches(spec.size_in), Inches(spec.size_in),
    )
    shape.line.fill.background()
    apply_radial_gradient(
        shape,
        inner_color=spec.color, outer_color=spec.color,
        inner_alpha=spec.alpha_pct, outer_alpha=0,
    )
    if spec.soft_edge_pt > 0:
        try:
            apply_soft_edge(shape, radius_pt=spec.soft_edge_pt)
        except Exception:
            pass


def _add_mesh(slide: Slide, spec: MeshSpec, w: float, h: float) -> None:
    rect = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(w), Inches(h),
    )
    rect.line.fill.background()
    stops = [(0, spec.colors[0], spec.alpha_pct),
             (50, spec.colors[1], spec.alpha_pct),
             (100, spec.colors[2], spec.alpha_pct)]
    apply_linear_gradient(rect, stops=stops, angle_deg=spec.angle_deg)


def _add_corner_glow(slide: Slide, spec: CornerGlowSpec, w: float, h: float) -> None:
    radius = max(w, h) * spec.radius_frac
    half = radius / 2
    cx, cy = {
        "tl": (0, 0), "tr": (w, 0), "bl": (0, h), "br": (w, h),
        "top": (w / 2, 0), "bottom": (w / 2, h),
    }[spec.corner]
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(cx - half), Inches(cy - half),
        Inches(radius), Inches(radius),
    )
    shape.line.fill.background()
    apply_radial_gradient(
        shape,
        inner_color=spec.color, outer_color=spec.color,
        inner_alpha=spec.alpha_pct, outer_alpha=0,
    )


def _add_edge_ribbon(slide: Slide, spec: EdgeRibbonSpec, w: float, h: float) -> None:
    if spec.edge == "left":
        x, y, sw, sh = spec.inset_in, 0, spec.width_in, h
    elif spec.edge == "right":
        x, y, sw, sh = w - spec.width_in - spec.inset_in, 0, spec.width_in, h
    elif spec.edge == "top":
        x, y, sw, sh = 0, spec.inset_in, w, spec.width_in
    else:
        x, y, sw, sh = 0, h - spec.width_in - spec.inset_in, w, spec.width_in
    rect = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(sw), Inches(sh),
    )
    rect.line.fill.background()
    apply_solid_fill(rect, spec.color, alpha=100)


def _add_geometric(slide: Slide, spec: GeometricSpec, w: float, h: float) -> None:
    x = spec.x * w if 0 <= spec.x <= 1.5 else spec.x
    y = spec.y * h if 0 <= spec.y <= 1.5 else spec.y
    half = spec.size_in / 2
    shape_map = {
        "triangle": MSO_SHAPE.ISOSCELES_TRIANGLE,
        "diamond": MSO_SHAPE.DIAMOND,
        "circle": MSO_SHAPE.OVAL,
        "pill": MSO_SHAPE.ROUNDED_RECTANGLE,
        "square": MSO_SHAPE.RECTANGLE,
    }
    sh = slide.shapes.add_shape(
        shape_map[spec.shape],
        Inches(x - half), Inches(y - half),
        Inches(spec.size_in), Inches(spec.size_in),
    )
    if spec.shape == "pill":
        sh.adjustments[0] = 0.5
    sh.line.fill.background()
    apply_solid_fill(sh, spec.color, alpha=spec.alpha_pct)
    if spec.rotation_deg:
        sh.rotation = spec.rotation_deg


def _add_scribble_dots(slide: Slide, spec: ScribbleDotsSpec, w: float, h: float) -> None:
    rng = random.Random(spec.seed)
    pad = 0.4
    radius_in = spec.radius_pt / 72.0
    for _ in range(spec.count):
        cx = rng.uniform(pad, w - pad)
        cy = rng.uniform(pad, h - pad)
        sh = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(cx - radius_in), Inches(cy - radius_in),
            Inches(radius_in * 2), Inches(radius_in * 2),
        )
        sh.line.fill.background()
        apply_solid_fill(sh, spec.color, alpha=spec.alpha_pct)


def _add_aurora_band(slide: Slide, spec: AuroraBandSpec, w: float, h: float) -> None:
    band_top = (h - spec.band_height_in) * spec.y_frac
    rect = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(-1.0), Inches(band_top),
        Inches(w + 2.0), Inches(spec.band_height_in),
    )
    rect.line.fill.background()
    n = len(spec.colors)
    stops = [(int(i * 100 / (n - 1)), c, 60) for i, c in enumerate(spec.colors)]
    apply_linear_gradient(rect, stops=stops, angle_deg=spec.angle_deg)
    rect.rotation = spec.angle_deg / 6  # subtle skew for organic feel


def _add_sticker(slide: Slide, spec: StickerSpec, w: float, h: float) -> None:
    x = spec.x * w if 0 <= spec.x <= 1.5 else spec.x
    y = spec.y * h if 0 <= spec.y <= 1.5 else spec.y
    rect = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y), Inches(spec.width_in), Inches(spec.height_in),
    )
    rect.adjustments[0] = 0.35
    rect.line.fill.background()
    apply_solid_fill(rect, spec.color, alpha=100)
    apply_drop_shadow(
        rect, color="000000", alpha_pct=25, blur_pt=8, distance_pt=2,
    )
    rect.rotation = spec.rotation_deg
    # Inline label
    tf = rect.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.10)
    tf.margin_right = Inches(0.10)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    p = tf.paragraphs[0]
    from pptx.enum.text import PP_ALIGN
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = spec.text
    r.font.bold = True
    r.font.size = Pt(11)
    r.font.color.rgb = _rgb(spec.text_color)


_RENDERERS = {
    "orb": _add_orb,
    "mesh": _add_mesh,
    "corner_glow": _add_corner_glow,
    "edge_ribbon": _add_edge_ribbon,
    "geometric": _add_geometric,
    "scribble_dots": _add_scribble_dots,
    "aurora_band": _add_aurora_band,
    "sticker": _add_sticker,
}


_SPEC_TYPES = {
    "orb": OrbSpec,
    "mesh": MeshSpec,
    "corner_glow": CornerGlowSpec,
    "edge_ribbon": EdgeRibbonSpec,
    "geometric": GeometricSpec,
    "scribble_dots": ScribbleDotsSpec,
    "aurora_band": AuroraBandSpec,
    "sticker": StickerSpec,
}


def parse_decoration(spec_dict: dict):
    """Validate a single decoration dict and return its parsed Pydantic model."""
    kind = spec_dict.get("kind")
    if kind not in _SPEC_TYPES:
        raise ValueError(f"Unknown decoration kind: {kind!r}. Valid: {list(_SPEC_TYPES)}")
    return _SPEC_TYPES[kind].model_validate(spec_dict)


def apply_decorations(slide: Slide, theme) -> None:
    """Render every decoration in `theme.decorations` onto `slide`.

    Decorations are laid down BEFORE content so they always sit underneath
    text. Each one is fail-safe in production, but in `MEGADECK_DEBUG=1`
    mode we re-raise so theme authors see typos & bad specs immediately.
    """
    import os
    debug = os.environ.get("MEGADECK_DEBUG") == "1"
    decos = getattr(theme, "decorations", None) or []
    if not decos:
        return
    w = theme.slide_width_in
    h = theme.slide_height_in
    for spec in decos:
        if isinstance(spec, dict):
            try:
                spec = parse_decoration(spec)
            except Exception:
                if debug:
                    raise
                continue
        kind = getattr(spec, "kind", None)
        if kind not in _RENDERERS:
            if debug:
                raise ValueError(f"Unknown decoration kind: {kind!r}")
            continue
        try:
            _RENDERERS[kind](slide, spec, w, h)
        except Exception:
            if debug:
                raise
