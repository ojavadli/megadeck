"""Slide compositions — *bold, unmistakable shape language*.

Compositions are the answer to "every slide looks the same with different
colours". Each composition decides *how* a slide is built up visually —
not just colour palette but actual shape vocabulary.

  * `typographic`  — pure type, no shapes other than negative space
  * `swiss`        — chunky horizontal accent banner + double hairline rule.
                     International Typographic Style by way of Müller-Brockmann.
  * `blueprint`    — visible graph paper grid + corner crop marks +
                     technical "MEGADECK / MM:SS" plate top-right.
  * `brutalist`    — heavy black title bar across the top, thick left rail,
                     loud bottom slab. Raw, unsentimental geometry.
  * `editorial`    — magazine-style strong vertical rule at 7/12 columns +
                     folio mark + drop-cap rule + horizontal kicker.
  * `photographic` — gradient-tinted right 45% photo zone with caption rail.
  * `grid`         — visible 12-column light grid + horizontal baseline.
  * `mono-grid`    — IBM-Plex-Mono technical grid with axis labels.
  * `bauhaus`      — primary-colour geometric blocks (red/yellow/blue) in
                     deliberate asymmetric placement.
  * `risograph`    — overlapping translucent shapes, one per accent colour.
  * `orbs`         — legacy default (kept for back-compat, opt-in only).

Compositions render AFTER the background fill but BEFORE template
content. They are non-destructive design layers.
"""
from __future__ import annotations

from typing import Callable, Dict, Optional

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.slide import Slide
from pptx.util import Inches, Pt

from megadeck.design_system.effects import apply_solid_fill
from megadeck.design_system.tokens import Theme


CompositionFn = Callable[[Slide, Theme], None]
COMPOSITIONS: Dict[str, CompositionFn] = {}


def register_composition(name: str):
    def deco(fn: CompositionFn) -> CompositionFn:
        COMPOSITIONS[name] = fn
        return fn
    return deco


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _hex_to_rgb(h: str) -> RGBColor:
    h = h.lstrip("#").upper().ljust(6, "0")[:6]
    return RGBColor.from_string(h)


def _add_rect(
    slide: Slide,
    *,
    left: float, top: float, width: float, height: float,
    color: RGBColor, no_line: bool = True,
):
    sh = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    apply_solid_fill(sh, color)
    if no_line:
        sh.line.fill.background()
    return sh


def _add_text(
    slide: Slide, text: str,
    *,
    left: float, top: float, w: float, h: float,
    color: RGBColor, font: str, size_pt: float,
    bold: bool = False, italic: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size_pt)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color


# ---------------------------------------------------------------------------
# typographic — pure space
# ---------------------------------------------------------------------------

@register_composition("typographic")
def _typographic(slide: Slide, theme: Theme) -> None:
    """Pure typography. The template's text is the only design element."""


# ---------------------------------------------------------------------------
# swiss — chunky accent banner top-left + double hairline rule bottom
# ---------------------------------------------------------------------------

@register_composition("swiss")
def _swiss(slide: Slide, theme: Theme) -> None:
    sw = theme.slide_width_in
    sh = theme.slide_height_in
    # Chunky accent banner top-left — was 0.45in × 0.10in, now 1.20in × 0.18in
    _add_rect(slide, left=theme.left_margin_in, top=0.50,
              width=1.20, height=0.18, color=theme.accent)
    # Double hairline rule near bottom — much bolder than before
    _add_rect(slide, left=theme.left_margin_in, top=sh - 0.55,
              width=sw - 2 * theme.left_margin_in, height=0.025,
              color=theme.title)
    _add_rect(slide, left=theme.left_margin_in, top=sh - 0.50,
              width=sw - 2 * theme.left_margin_in, height=0.012,
              color=theme.hairline)
    # Right-edge column rule (Müller-Brockmann signature)
    _add_rect(slide, left=sw - 0.40, top=0.50,
              width=0.025, height=sh - 1.0,
              color=theme.title)


# ---------------------------------------------------------------------------
# blueprint — graph-paper visible + corner crop marks + technical plate
# ---------------------------------------------------------------------------

@register_composition("blueprint")
def _blueprint(slide: Slide, theme: Theme) -> None:
    sw = theme.slide_width_in
    sh = theme.slide_height_in
    grid_color = theme.hairline

    # Major grid every 1.0in (visible)
    for i in range(0, int(sw) + 1):
        _add_rect(slide, left=float(i) - 0.005, top=0.0,
                  width=0.012, height=sh, color=grid_color)
    for i in range(0, int(sh) + 1):
        _add_rect(slide, left=0.0, top=float(i) - 0.005,
                  width=sw, height=0.012, color=grid_color)
    # Minor grid every 0.25in (very faint)
    light_grid = theme.muted
    step = 0.25
    x = step
    while x < sw:
        _add_rect(slide, left=x, top=0.0, width=0.005, height=sh, color=light_grid)
        x += step
    y = step
    while y < sh:
        _add_rect(slide, left=0.0, top=y, width=sw, height=0.005, color=light_grid)
        y += step
    # Corner crop marks — bigger and bolder
    L = 0.40
    crop = theme.title
    for cx, cy in (
        (0.20, 0.20), (sw - 0.60, 0.20),
        (0.20, sh - 0.60), (sw - 0.60, sh - 0.60),
    ):
        _add_rect(slide, left=cx, top=cy, width=L, height=0.020, color=crop)
        _add_rect(slide, left=cx, top=cy, width=0.020, height=L, color=crop)
    # Technical plate top-right
    _add_text(
        slide, "DOC.001 / R0",
        left=sw - 1.80, top=0.30, w=1.50, h=0.30,
        color=theme.muted, font="IBM Plex Mono",
        size_pt=9, align=PP_ALIGN.RIGHT,
    )


# ---------------------------------------------------------------------------
# brutalist — heavy slabs everywhere
# ---------------------------------------------------------------------------

@register_composition("brutalist")
def _brutalist(slide: Slide, theme: Theme) -> None:
    sw = theme.slide_width_in
    sh = theme.slide_height_in
    # Heavy left rail — 0.40in (was 0.20)
    _add_rect(slide, left=0.0, top=0.0, width=0.40, height=sh, color=theme.title)
    # Heavy bottom slab — 0.32in (was 0.18)
    _add_rect(slide, left=0.0, top=sh - 0.32, width=sw, height=0.32, color=theme.title)
    # Top-right loud accent slab
    _add_rect(slide, left=sw - 2.40, top=0.40, width=2.40, height=0.65, color=theme.accent)
    # White "01" inside the slab — brutalist numbering convention
    _add_text(
        slide, "01",
        left=sw - 2.20, top=0.45, w=0.80, h=0.55,
        color=_hex_to_rgb("FFFFFF"),
        font=theme.font_display,
        size_pt=44, bold=True,
    )


# ---------------------------------------------------------------------------
# editorial — magazine asymmetric grid + drop-cap rule
# ---------------------------------------------------------------------------

@register_composition("editorial")
def _editorial(slide: Slide, theme: Theme) -> None:
    sw = theme.slide_width_in
    sh = theme.slide_height_in
    # Strong vertical rule at the 7/12 column position
    rule_x = sw * (7 / 12)
    _add_rect(slide, left=rule_x, top=0.65, width=0.020, height=sh - 1.30,
              color=theme.title)
    # Top horizontal kicker rule above title area
    _add_rect(slide, left=theme.left_margin_in, top=1.06,
              width=2.50, height=0.025, color=theme.accent)
    # Folio mark top-right (issue / page)
    _add_text(
        slide, "ISSUE 01",
        left=sw - 1.50, top=0.40, w=1.20, h=0.30,
        color=theme.muted, font=theme.font_body,
        size_pt=9, italic=True, align=PP_ALIGN.RIGHT,
    )
    # Bottom-edge editorial signature line
    _add_rect(slide, left=theme.left_margin_in, top=sh - 0.60,
              width=sw - 2 * theme.left_margin_in, height=0.012,
              color=theme.hairline)
    _add_text(
        slide, "—  Megadeck Editorial",
        left=theme.left_margin_in, top=sh - 0.45, w=4.0, h=0.30,
        color=theme.muted, font=theme.font_body,
        size_pt=9, italic=True,
    )


# ---------------------------------------------------------------------------
# photographic — strong photo zone right
# ---------------------------------------------------------------------------

@register_composition("photographic")
def _photographic(slide: Slide, theme: Theme) -> None:
    sw = theme.slide_width_in
    sh = theme.slide_height_in
    photo_w = sw * 0.45
    photo_l = sw - photo_w
    _add_rect(slide, left=photo_l, top=0.0, width=photo_w, height=sh,
              color=theme.accent_lt)
    # Caption rail at the bottom
    _add_rect(slide, left=photo_l, top=sh - 0.50,
              width=photo_w, height=0.50,
              color=theme.title)
    _add_text(
        slide, "FIG. 01  /  MEGADECK",
        left=photo_l + 0.30, top=sh - 0.42,
        w=photo_w - 0.60, h=0.30,
        color=_hex_to_rgb("FFFFFF"),
        font=theme.font_body, size_pt=10,
    )


# ---------------------------------------------------------------------------
# grid — visible 12-col grid + baseline
# ---------------------------------------------------------------------------

@register_composition("grid")
def _grid(slide: Slide, theme: Theme) -> None:
    sw = theme.slide_width_in
    sh = theme.slide_height_in
    margin = theme.left_margin_in
    cells = 12
    inner_w = sw - 2 * margin
    col_w = inner_w / cells
    # Visible column rules — much bolder than before (0.014 vs 0.005)
    for i in range(cells + 1):
        x = margin + i * col_w
        _add_rect(slide, left=x, top=0.50,
                  width=0.014, height=sh - 1.00,
                  color=theme.hairline)
    # Strong horizontal baseline rule top + bottom
    _add_rect(slide, left=margin, top=0.50, width=inner_w, height=0.025,
              color=theme.title)
    _add_rect(slide, left=margin, top=sh - 0.55, width=inner_w, height=0.025,
              color=theme.title)


# ---------------------------------------------------------------------------
# mono-grid — IBM Plex Mono technical grid with axis labels
# ---------------------------------------------------------------------------

@register_composition("mono-grid")
def _mono_grid(slide: Slide, theme: Theme) -> None:
    sw = theme.slide_width_in
    sh = theme.slide_height_in
    # Mono ruler labels every 1.0in across the top
    for i in range(0, int(sw) + 1):
        _add_text(
            slide, f"{i:02d}",
            left=float(i) - 0.10, top=0.20, w=0.40, h=0.20,
            color=theme.muted, font="IBM Plex Mono", size_pt=8,
        )
        _add_rect(slide, left=float(i), top=0.40, width=0.012, height=0.10,
                  color=theme.muted)
    # Top hairline
    _add_rect(slide, left=0.0, top=0.50, width=sw, height=0.012,
              color=theme.hairline)
    # Bottom plate
    _add_rect(slide, left=0.0, top=sh - 0.40, width=sw, height=0.012,
              color=theme.hairline)
    _add_text(
        slide, "DECK / MEGADECK / VERSION 0.4",
        left=theme.left_margin_in, top=sh - 0.30, w=4.0, h=0.20,
        color=theme.muted, font="IBM Plex Mono", size_pt=8,
    )


# ---------------------------------------------------------------------------
# bauhaus — primary-colour geometric blocks
# ---------------------------------------------------------------------------

@register_composition("bauhaus")
def _bauhaus(slide: Slide, theme: Theme) -> None:
    sw = theme.slide_width_in
    sh = theme.slide_height_in
    # Big yellow square top-right — Bauhaus primaries (we cheat and use accent
    # variants if the theme doesn't have true primaries).
    yellow = _hex_to_rgb("FFD500")
    red = _hex_to_rgb("E53935")
    blue = _hex_to_rgb("1565C0")
    # Bottom-right blue triangle (using a right triangle shape)
    triangle = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_TRIANGLE,
        Inches(sw - 1.80), Inches(sh - 1.80),
        Inches(1.80), Inches(1.80),
    )
    apply_solid_fill(triangle, blue)
    triangle.line.fill.background()
    # Top-right yellow square
    _add_rect(slide, left=sw - 1.30, top=0.30, width=1.0, height=1.0,
              color=yellow)
    # Mid-bottom red circle
    circ = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(0.40), Inches(sh - 1.20),
        Inches(0.70), Inches(0.70),
    )
    apply_solid_fill(circ, red)
    circ.line.fill.background()


# ---------------------------------------------------------------------------
# risograph — overlapping translucent shapes
# ---------------------------------------------------------------------------

@register_composition("risograph")
def _risograph(slide: Slide, theme: Theme) -> None:
    """Overlapping circles in two ink colours, low-fi print aesthetic."""
    from megadeck.design_system.effects import apply_solid_fill
    sw = theme.slide_width_in
    sh = theme.slide_height_in
    # Two overlapping circles bottom-right — pinkish + teal-ish from theme accents
    c1 = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(sw - 3.00), Inches(sh - 2.50),
        Inches(2.40), Inches(2.40),
    )
    apply_solid_fill(c1, theme.accent, alpha=70)
    c1.line.fill.background()
    c2 = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(sw - 1.80), Inches(sh - 2.00),
        Inches(2.00), Inches(2.00),
    )
    apply_solid_fill(c2, theme.accent_lt, alpha=80)
    c2.line.fill.background()


# ---------------------------------------------------------------------------
# futurist — Apple-style minimal HUD: thin grid + corner brackets + status pill
# ---------------------------------------------------------------------------

@register_composition("futurist")
def _futurist(slide: Slide, theme: Theme) -> None:
    """Futuristic Apple-HUD aesthetic: ambient grid, corner brackets, glass pill."""
    from megadeck.design_system.effects import apply_solid_fill, apply_radial_gradient
    sw = theme.slide_width_in
    sh = theme.slide_height_in

    # Ambient large radial glow top-right (subtle accent)
    glow = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(sw - 5.50), Inches(-2.80),
        Inches(7.50), Inches(7.50),
    )
    glow.line.fill.background()
    try:
        apply_radial_gradient(
            glow,
            inner_color=theme.accent, outer_color=theme.accent,
            inner_alpha=14, outer_alpha=0,
        )
    except Exception:
        pass

    # Ambient large radial glow bottom-left (deeper accent)
    glow2 = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(-3.00), Inches(sh - 4.00),
        Inches(7.00), Inches(7.00),
    )
    glow2.line.fill.background()
    try:
        apply_radial_gradient(
            glow2,
            inner_color=theme.accent_dk, outer_color=theme.accent_dk,
            inner_alpha=10, outer_alpha=0,
        )
    except Exception:
        pass

    # Thin grid lines — only every 2 inches (very subtle)
    for x in (sw / 4, sw / 2, 3 * sw / 4):
        _add_rect(slide, left=x, top=0.50, width=0.008, height=sh - 1.00,
                  color=theme.hairline)
    # Horizontal baseline rules top + bottom (thin)
    _add_rect(slide, left=0.50, top=0.50, width=sw - 1.00, height=0.012,
              color=theme.hairline)
    _add_rect(slide, left=0.50, top=sh - 0.55, width=sw - 1.00, height=0.012,
              color=theme.hairline)

    # Corner L-brackets (HUD aesthetic) — 4 corners
    L = 0.32
    bracket = theme.accent
    for cx, cy in (
        (0.30, 0.30), (sw - 0.62, 0.30),
        (0.30, sh - 0.42), (sw - 0.62, sh - 0.42),
    ):
        _add_rect(slide, left=cx, top=cy, width=L, height=0.018, color=bracket)
        _add_rect(slide, left=cx, top=cy, width=0.018, height=L, color=bracket)

    # Status pill top-right (glass effect simulated with thin border + low-alpha fill)
    pill_w, pill_h = 1.40, 0.32
    pill_x = sw - pill_w - 0.40
    pill_y = 0.30
    pill = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(pill_x), Inches(pill_y),
        Inches(pill_w), Inches(pill_h),
    )
    pill.adjustments[0] = 0.50
    apply_solid_fill(pill, theme.surface, alpha=8)
    try:
        pill.line.color.rgb = theme.hairline
        pill.line.width = Pt(0.5)
    except Exception:
        pass
    # Status dot
    dot = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(pill_x + 0.12), Inches(pill_y + (pill_h - 0.10) / 2),
        Inches(0.10), Inches(0.10),
    )
    apply_solid_fill(dot, theme.accent)
    dot.line.fill.background()
    _add_text(
        slide, "MEGADECK / 2026",
        left=pill_x + 0.30, top=pill_y + 0.05,
        w=pill_w - 0.30, h=0.22,
        color=theme.body, font="SF Pro Text", size_pt=8,
    )


# ---------------------------------------------------------------------------
# orbs — legacy default (still available, opt-in only)
# ---------------------------------------------------------------------------

@register_composition("orbs")
def _orbs(slide: Slide, theme: Theme) -> None:
    """Legacy default: glowing orbs. Kept for back-compat."""
    from megadeck.design_system.decorations import apply_decorations
    apply_decorations(slide, theme)


# ---------------------------------------------------------------------------
# Public dispatch
# ---------------------------------------------------------------------------

def apply_composition(slide: Slide, theme: Theme, *, override: Optional[str] = None) -> str:
    """Apply a composition. Returns the composition name actually rendered."""
    requested = override or getattr(theme, "composition", None) or "typographic"
    fn = COMPOSITIONS.get(requested)
    if fn is None:
        fn = COMPOSITIONS["typographic"]
        requested = "typographic"
    fn(slide, theme)
    return requested


# Default rotation order — used by rhythm orchestrator. Adjacent slides
# get visually different compositions.
DEFAULT_ROTATION = [
    "typographic", "swiss", "editorial", "blueprint",
    "grid", "brutalist", "photographic", "mono-grid",
    "bauhaus", "risograph",
]


def composition_at_index(idx: int, *, rotation=DEFAULT_ROTATION) -> str:
    return rotation[idx % len(rotation)]


__all__ = [
    "COMPOSITIONS", "register_composition",
    "apply_composition", "composition_at_index", "DEFAULT_ROTATION",
]
