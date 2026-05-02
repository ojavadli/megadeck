"""Layout variants for `numbered_list`.

Each variant takes the same `NumberedListSlide` data and renders a *structurally
different* composition. This is what gives a deck visual variety without
re-authoring content — choose the variant per slide, or set a theme-level
default.

Variants
--------

* `default`  — classic outlined-numerals + vertical accent bar (existing).
                Lives in `numbered_list.py`.

* `split`    — numbers anchored to the right edge in a tall column, item
                text fills the left 70% in a single tall card. A more
                "magazine" feel, asymmetric.

* `cards`    — each item lives inside its own rounded-corner card with a
                badge number top-left. Reads as a board / dashboard.

* `timeline` — vertical timeline: a single accent line on the left,
                round number nodes anchored to it, item text fans right.
                Great for sequential / process content.
"""
from __future__ import annotations

from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.slide import Slide
from pptx.util import Inches, Pt

from megadeck.core.schemas import NumberedListSlide
from megadeck.design_system.primitives import (
    add_corner_dotgrid,
    add_eyebrow,
    add_oval,
    add_page_chrome,
    add_round_rect,
    add_text,
    add_themed_card,
    add_v_line,
    fit_title,
    measure_title_height,
    set_slide_bg,
)
from megadeck.design_system.tokens import Theme
from megadeck.design_system.variants import register_variant


# ---------------------------------------------------------------------------
# split — numbers in a tall right column, text on the left
# ---------------------------------------------------------------------------

@register_variant("numbered_list", "split")
def render_numbered_list_split(
    slide: Slide,
    data: NumberedListSlide,
    theme: Theme,
    *,
    page_n: int,
    page_total: int,
    section_label: str | None = None,
) -> None:
    set_slide_bg(slide, color=theme.bg, theme=theme)
    add_corner_dotgrid(slide, theme)
    add_eyebrow(slide, text=data.eyebrow.upper(), theme=theme)

    LEFT = theme.left_margin_in
    CONTENT_W = theme.content_width_in
    SLIDE_W = theme.slide_width_in

    title_pt = fit_title(data.title, max_pt=theme.type_scale.h2, width_in=CONTENT_W)
    title_h = measure_title_height(data.title, size_pt=title_pt, width_in=CONTENT_W)
    add_text(
        slide,
        left=LEFT, top=1.20, width=CONTENT_W, height=title_h,
        text=data.title,
        font=theme.font_display,
        size_pt=title_pt,
        color=theme.title,
        bold=True,
        line_spacing=1.05,
    )

    # The split: text on the left 65%, numbers on the right 25% column.
    body_top = 1.20 + title_h + 0.40
    bottom = 6.95
    available = max(0.6, bottom - body_top)
    n = len(data.items)
    gap = 0.24 if n <= 4 else 0.16
    item_h = max(0.55, (available - (n - 1) * gap) / max(n, 1))

    TEXT_X = LEFT
    TEXT_W = CONTENT_W * 0.66
    NUM_X = LEFT + CONTENT_W * 0.74
    NUM_W = CONTENT_W * 0.26

    # Tall vertical hairline between text + number columns
    add_v_line(
        slide,
        left=NUM_X - 0.18, top=body_top,
        height=available, color=theme.hairline, width=0.012,
    )

    head_pt = 22 if n <= 4 else 19 if n == 5 else 17
    tail_pt = head_pt - 3
    num_pt = 64 if n <= 3 else 52 if n <= 4 else 44 if n <= 5 else 38

    for i, item in enumerate(data.items):
        y = body_top + i * (item_h + gap)
        # Right column: big outlined number, right-aligned
        add_text(
            slide,
            left=NUM_X, top=y, width=NUM_W, height=item_h,
            text=f"{i+1:02d}",
            font=theme.font_display,
            size_pt=num_pt,
            color=theme.accent_lt,
            bold=True,
            align=PP_ALIGN.RIGHT,
            anchor=MSO_ANCHOR.MIDDLE,
            line_spacing=0.95,
        )
        # Left column: head + tail
        tb = slide.shapes.add_textbox(
            Inches(TEXT_X), Inches(y), Inches(TEXT_W), Inches(item_h),
        )
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.04)
        try:
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        except Exception:
            pass
        p = tf.paragraphs[0]
        try:
            p.line_spacing = 1.20
        except Exception:
            pass
        r1 = p.add_run()
        r1.text = item.head
        r1.font.name = theme.font_display
        r1.font.size = Pt(head_pt)
        r1.font.bold = True
        r1.font.color.rgb = theme.title
        if item.tail:
            p2 = tf.add_paragraph()
            try:
                p2.line_spacing = 1.30
            except Exception:
                pass
            r2 = p2.add_run()
            r2.text = item.tail
            r2.font.name = theme.font_body
            r2.font.size = Pt(tail_pt)
            r2.font.color.rgb = theme.body

    add_page_chrome(
        slide, theme=theme,
        page_n=page_n, page_total=page_total,
        section_label=section_label,
    )


# ---------------------------------------------------------------------------
# cards — each item in its own card
# ---------------------------------------------------------------------------

@register_variant("numbered_list", "cards")
def render_numbered_list_cards(
    slide: Slide,
    data: NumberedListSlide,
    theme: Theme,
    *,
    page_n: int,
    page_total: int,
    section_label: str | None = None,
) -> None:
    set_slide_bg(slide, color=theme.bg, theme=theme)
    add_corner_dotgrid(slide, theme)
    add_eyebrow(slide, text=data.eyebrow.upper(), theme=theme)

    LEFT = theme.left_margin_in
    CONTENT_W = theme.content_width_in

    title_pt = fit_title(data.title, max_pt=theme.type_scale.h2, width_in=CONTENT_W)
    title_h = measure_title_height(data.title, size_pt=title_pt, width_in=CONTENT_W)
    add_text(
        slide,
        left=LEFT, top=1.20, width=CONTENT_W, height=title_h,
        text=data.title,
        font=theme.font_display,
        size_pt=title_pt,
        color=theme.title,
        bold=True,
        line_spacing=1.05,
    )

    grid_top = 1.20 + title_h + 0.40
    bottom = 6.95
    available = max(0.8, bottom - grid_top)
    n = len(data.items)
    # 2-column grid for 4-6 items, single column for ≤3
    if n <= 3:
        cols = 1
    elif n <= 6:
        cols = 2
    else:
        cols = 3
    rows = (n + cols - 1) // cols
    col_gap = 0.24
    row_gap = 0.20
    col_w = (CONTENT_W - (cols - 1) * col_gap) / cols
    row_h = max(1.30, (available - (rows - 1) * row_gap) / max(rows, 1))

    head_pt = 18 if n <= 4 else 16
    tail_pt = head_pt - 3

    for i, item in enumerate(data.items):
        col = i % cols
        row = i // cols
        x = LEFT + col * (col_w + col_gap)
        y = grid_top + row * (row_h + row_gap)
        add_themed_card(
            slide, theme,
            left=x, top=y, width=col_w, height=row_h,
            adjust=0.05,
        )
        # Big number badge top-left
        badge_size = 0.55
        add_oval(
            slide,
            left=x + 0.30, top=y + 0.30,
            size=badge_size,
            fill=theme.accent,
        )
        add_text(
            slide,
            left=x + 0.30, top=y + 0.30, width=badge_size, height=badge_size,
            text=f"{i+1}",
            font=theme.font_display,
            size_pt=22,
            color=theme.inverse,
            bold=True,
            align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        # Head
        add_text(
            slide,
            left=x + 0.30 + badge_size + 0.18, top=y + 0.30,
            width=col_w - 0.78 - badge_size, height=badge_size,
            text=item.head,
            font=theme.font_display,
            size_pt=head_pt,
            color=theme.title,
            bold=True,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        # Tail
        if item.tail:
            add_text(
                slide,
                left=x + 0.30, top=y + 0.30 + badge_size + 0.20,
                width=col_w - 0.60, height=row_h - 0.50 - badge_size,
                text=item.tail,
                font=theme.font_body,
                size_pt=tail_pt,
                color=theme.body,
                line_spacing=1.30,
            )

    add_page_chrome(
        slide, theme=theme,
        page_n=page_n, page_total=page_total,
        section_label=section_label,
    )


# ---------------------------------------------------------------------------
# timeline — vertical timeline with number nodes
# ---------------------------------------------------------------------------

@register_variant("numbered_list", "timeline")
def render_numbered_list_timeline(
    slide: Slide,
    data: NumberedListSlide,
    theme: Theme,
    *,
    page_n: int,
    page_total: int,
    section_label: str | None = None,
) -> None:
    set_slide_bg(slide, color=theme.bg, theme=theme)
    add_corner_dotgrid(slide, theme)
    add_eyebrow(slide, text=data.eyebrow.upper(), theme=theme)

    LEFT = theme.left_margin_in
    CONTENT_W = theme.content_width_in

    title_pt = fit_title(data.title, max_pt=theme.type_scale.h2, width_in=CONTENT_W)
    title_h = measure_title_height(data.title, size_pt=title_pt, width_in=CONTENT_W)
    add_text(
        slide,
        left=LEFT, top=1.20, width=CONTENT_W, height=title_h,
        text=data.title,
        font=theme.font_display,
        size_pt=title_pt,
        color=theme.title,
        bold=True,
        line_spacing=1.05,
    )

    body_top = 1.20 + title_h + 0.40
    bottom = 6.95
    available = max(0.6, bottom - body_top)
    n = len(data.items)
    gap = 0.20
    # Floor item_h at 0.55in — items shouldn't get squeezed below that.
    raw_h = (available - (n - 1) * gap) / max(n, 1)
    item_h = max(0.55, raw_h)
    # If we hit the floor, recompute the consumed height so the rail matches
    # the visible items rather than the original `available`.
    consumed_h = n * item_h + (n - 1) * gap

    # Vertical timeline rail
    NODE_X = LEFT + 0.40
    add_v_line(
        slide,
        left=NODE_X, top=body_top,
        height=min(available, consumed_h), color=theme.accent_lt, width=0.04,
    )

    NODE_SIZE = 0.70
    TEXT_X = NODE_X + NODE_SIZE + 0.40
    TEXT_W = CONTENT_W - (TEXT_X - LEFT)

    head_pt = 20 if n <= 4 else 17
    tail_pt = head_pt - 3

    for i, item in enumerate(data.items):
        y = body_top + i * (item_h + gap)
        node_y = y + (item_h - NODE_SIZE) / 2
        # Number node
        add_oval(
            slide,
            left=NODE_X - NODE_SIZE / 2 + 0.02, top=node_y,
            size=NODE_SIZE,
            fill=theme.accent,
        )
        add_text(
            slide,
            left=NODE_X - NODE_SIZE / 2 + 0.02, top=node_y,
            width=NODE_SIZE, height=NODE_SIZE,
            text=f"{i+1}",
            font=theme.font_display,
            size_pt=22,
            color=theme.inverse,
            bold=True,
            align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        # Item text
        tb = slide.shapes.add_textbox(
            Inches(TEXT_X), Inches(y), Inches(TEXT_W), Inches(item_h),
        )
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.04)
        try:
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        except Exception:
            pass
        p = tf.paragraphs[0]
        try:
            p.line_spacing = 1.20
        except Exception:
            pass
        r1 = p.add_run()
        r1.text = item.head
        r1.font.name = theme.font_display
        r1.font.size = Pt(head_pt)
        r1.font.bold = True
        r1.font.color.rgb = theme.title
        if item.tail:
            r2 = p.add_run()
            r2.text = "  " + item.tail
            r2.font.name = theme.font_body
            r2.font.size = Pt(tail_pt)
            r2.font.color.rgb = theme.body

    add_page_chrome(
        slide, theme=theme,
        page_n=page_n, page_total=page_total,
        section_label=section_label,
    )
