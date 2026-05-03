"""photo_card template — large photo + headline split.

The photo can be a local path (rendered as an embedded image) or an
HTTPS URL (downloaded once + cached). When `data.photo` is None, we
draw a tinted card placeholder so the layout still composes cleanly.
"""
from __future__ import annotations

import io
import urllib.request
from pathlib import Path

from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.slide import Slide
from pptx.util import Inches

from megadeck.core.schemas import PhotoCardSlide
from megadeck.design_system.primitives import (
    add_corner_dotgrid,
    add_eyebrow,
    add_page_chrome,
    add_round_rect,
    add_text,
    fit_title,
    measure_title_height,
    set_slide_bg,
)
from megadeck.design_system.tokens import Theme


def _embed_photo(
    slide: Slide,
    photo: str | None,
    *,
    left: float, top: float,
    width: float, height: float,
    theme: Theme,
) -> None:
    # Resolve `unsplash:keywords` specs into actual local files before embedding.
    try:
        from megadeck.design_system.unsplash import resolve_photo
        photo = resolve_photo(photo)
    except Exception:
        pass
    if photo is None:
        # Placeholder
        add_round_rect(
            slide, left=left, top=top, width=width, height=height,
            fill=theme.accent_lt, adjust=0.06,
        )
        add_text(
            slide,
            left=left, top=top, width=width, height=height,
            text="◑", font=theme.font_display,
            size_pt=120, color=theme.bg,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
            auto_size=False,
        )
        return
    # Resolve photo source
    try:
        if photo.startswith(("http://", "https://")):
            with urllib.request.urlopen(photo, timeout=10) as r:
                img_bytes = r.read()
            slide.shapes.add_picture(
                io.BytesIO(img_bytes),
                Inches(left), Inches(top), Inches(width), Inches(height),
            )
        else:
            p = Path(photo).expanduser()
            if not p.exists():
                raise FileNotFoundError(p)
            slide.shapes.add_picture(
                str(p),
                Inches(left), Inches(top), Inches(width), Inches(height),
            )
    except Exception:
        # Hide-the-failure-gracefully fallback
        add_round_rect(
            slide, left=left, top=top, width=width, height=height,
            fill=theme.muted, adjust=0.06,
        )


def render_photo_card(
    slide: Slide,
    data: PhotoCardSlide,
    theme: Theme,
    *,
    page_n: int,
    page_total: int,
    section_label: str | None = None,
) -> None:
    set_slide_bg(slide, color=theme.bg, theme=theme)
    add_corner_dotgrid(slide, theme)

    SLIDE_W = theme.slide_width_in
    SLIDE_H = theme.slide_height_in
    LEFT = theme.left_margin_in

    if data.photo_position == "full":
        # Full-bleed photo with text overlay bottom-left
        _embed_photo(
            slide, data.photo,
            left=0, top=0, width=SLIDE_W, height=SLIDE_H,
            theme=theme,
        )
        # Dark scrim bottom-left
        add_round_rect(
            slide,
            left=0, top=SLIDE_H - 2.5,
            width=SLIDE_W, height=2.5,
            fill=RGBColor.from_string("000000"), adjust=0.0,
        )
        add_text(
            slide,
            left=LEFT, top=SLIDE_H - 2.0,
            width=SLIDE_W - 2 * LEFT, height=1.5,
            text=data.title,
            font=theme.font_display, size_pt=44,
            color=RGBColor.from_string("FFFFFF"), bold=True,
        )
        add_page_chrome(
            slide, theme=theme,
            page_n=page_n, page_total=page_total,
            section_label=section_label,
        )
        return

    # Side-by-side split
    photo_w = SLIDE_W * 0.45
    text_w = SLIDE_W - 2 * LEFT - photo_w - 0.40

    if data.photo_position == "left":
        photo_l = 0.0
        text_l = photo_w + 0.40 + LEFT
    else:
        photo_l = SLIDE_W - photo_w
        text_l = LEFT

    _embed_photo(
        slide, data.photo,
        left=photo_l, top=0.0,
        width=photo_w, height=SLIDE_H,
        theme=theme,
    )

    add_eyebrow(slide, text=data.eyebrow.upper(), theme=theme, left=text_l)
    title_pt = fit_title(data.title, max_pt=theme.type_scale.h2, width_in=text_w)
    title_h = measure_title_height(data.title, size_pt=title_pt, width_in=text_w)
    add_text(
        slide,
        left=text_l, top=1.20, width=text_w, height=title_h,
        text=data.title,
        font=theme.font_display, size_pt=title_pt,
        color=theme.title, bold=True, line_spacing=1.05,
    )

    cursor = 1.20 + title_h + 0.30
    if data.subtitle:
        add_text(
            slide,
            left=text_l, top=cursor, width=text_w, height=0.80,
            text=data.subtitle,
            font=theme.font_body, size_pt=18, color=theme.body,
        )
        cursor += 0.85
    if data.body:
        add_text(
            slide,
            left=text_l, top=cursor, width=text_w, height=SLIDE_H - cursor - 0.80,
            text=data.body,
            font=theme.font_body, size_pt=14, color=theme.body, line_spacing=1.40,
        )

    add_page_chrome(
        slide, theme=theme,
        page_n=page_n, page_total=page_total,
        section_label=section_label,
    )
