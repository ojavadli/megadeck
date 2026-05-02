"""Editorial split — magazine-style 50/50 layout: text on one side,
decorative orb composition on the other. Asymmetric, breathing, distinct
from any list-style slide.

Content side gets eyebrow + title + body paragraph.
Decoration side gets a stacked composition: large background orb, smaller
foreground orb in a complementary hue, a thin accent stripe.
"""
from __future__ import annotations

from pptx.enum.shapes import MSO_SHAPE
from pptx.slide import Slide
from pptx.util import Inches

from megadeck.core.schemas import EditorialSplitSlide
from megadeck.design_system.effects import (
    apply_radial_gradient,
    apply_solid_fill,
    apply_soft_edge,
)
from megadeck.design_system.primitives import (
    add_eyebrow,
    add_page_chrome,
    add_text,
    add_v_line,
    fit_title,
    measure_title_height,
    set_slide_bg,
)
from megadeck.design_system.tokens import Theme


def _add_decoration_orbs(
    slide: Slide, theme: Theme, *, x: float, y: float, w: float, h: float,
    accent: str,
) -> None:
    """Drop a stacked orb composition centered on (x+w/2, y+h/2)."""
    cx, cy = x + w / 2, y + h / 2
    # Big back orb — accent colour, very soft
    big = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(cx - 2.4), Inches(cy - 2.4),
        Inches(4.8), Inches(4.8),
    )
    big.line.fill.background()
    apply_radial_gradient(
        big, inner_color=accent, outer_color=accent,
        inner_alpha=70, outer_alpha=0,
    )
    try:
        apply_soft_edge(big, radius_pt=32)
    except Exception:
        pass
    # Mid orb — accent_lt
    mid = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(cx - 1.4), Inches(cy - 1.4),
        Inches(2.8), Inches(2.8),
    )
    mid.line.fill.background()
    apply_radial_gradient(
        mid, inner_color=str(theme.accent_lt), outer_color=str(theme.accent_lt),
        inner_alpha=85, outer_alpha=0,
    )
    try:
        apply_soft_edge(mid, radius_pt=24)
    except Exception:
        pass
    # Foreground sphere — solid disc with soft edge
    fg = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(cx - 0.6), Inches(cy - 0.6),
        Inches(1.2), Inches(1.2),
    )
    fg.line.fill.background()
    apply_solid_fill(fg, accent)


def render_editorial_split(
    slide: Slide,
    data: EditorialSplitSlide,
    theme: Theme,
    *,
    page_n: int,
    page_total: int,
    section_label: str | None = None,
) -> None:
    set_slide_bg(slide, color=theme.bg, theme=theme)

    LEFT = theme.left_margin_in
    SLIDE_W = theme.slide_width_in
    CONTENT_W = theme.content_width_in
    HALF_W = (CONTENT_W - 0.40) / 2  # leave a 0.4-in gutter

    accent_hex = data.decoration_color or f"#{str(theme.accent)}"
    if not accent_hex.startswith("#"):
        accent_hex = "#" + accent_hex

    # Decide which half holds text
    if data.side == "left":
        text_x = LEFT
        deco_x = LEFT + HALF_W + 0.40
    else:
        text_x = LEFT + HALF_W + 0.40
        deco_x = LEFT

    add_eyebrow(slide, text=data.eyebrow.upper(), theme=theme,
                left=text_x, top=0.85, width=HALF_W)

    # Title — fit to half-width. auto_size=False because in OOXML
    # SHAPE_TO_FIT_TEXT can grow horizontally past the canvas; we use the
    # measured height directly with an extra safety pad instead.
    title_pt = fit_title(data.title, max_pt=theme.type_scale.h2, width_in=HALF_W)
    title_h = measure_title_height(data.title, size_pt=title_pt, width_in=HALF_W)
    add_text(
        slide,
        left=text_x, top=1.30, width=HALF_W, height=title_h + 0.30,
        text=data.title,
        font=theme.font_display,
        size_pt=title_pt,
        color=theme.title,
        bold=True,
        line_spacing=1.05,
        auto_size=False,
    )

    body_top = 1.30 + title_h + 0.30
    add_text(
        slide,
        left=text_x, top=body_top, width=HALF_W,
        height=theme.slide_height_in - body_top - 1.0,
        text=data.body,
        font=theme.font_body,
        size_pt=theme.type_scale.body_large,
        color=theme.body,
        line_spacing=1.40,
    )

    # Vertical hairline divider for editorial feel
    add_v_line(
        slide,
        left=LEFT + HALF_W + 0.20, top=1.30,
        height=theme.slide_height_in - 2.60,
        color=theme.hairline, width=0.012,
    )

    # Editorial decoration: a typographic plate (large pull-quote glyph
    # + accent rule + folio mark). NO MORE ORBS — this composition is now
    # purely typographic, like a magazine layout.
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    add_text(
        slide,
        left=deco_x, top=1.40, width=HALF_W, height=2.40,
        text="❝",
        font=theme.font_display,
        size_pt=320,
        color=theme.accent_lt,
        align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE,
        auto_size=False,
    )
    # Accent rule under the glyph
    add_v_line(
        slide,
        left=deco_x + HALF_W * 0.30, top=4.10,
        height=0.05, color=theme.accent, width=HALF_W * 0.40,
    )
    # Folio mark — small kicker text
    add_text(
        slide,
        left=deco_x, top=4.40, width=HALF_W, height=0.50,
        text=data.eyebrow.upper(),
        font=theme.font_body,
        size_pt=10,
        color=theme.muted,
        align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE,
        auto_size=False,
    )

    add_page_chrome(
        slide, theme=theme,
        page_n=page_n, page_total=page_total,
        section_label=section_label,
    )
