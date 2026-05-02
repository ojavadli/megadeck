"""Low-level shape helpers used by every template.

These wrap the verbose `python-pptx` API into a small, expressive vocabulary
(`add_text`, `add_card`, `add_v_line`, `add_dot_grid`, …) so that template
files stay short and readable.
"""
from __future__ import annotations

import math
from typing import Iterable, Optional, Sequence, Tuple

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.slide import Slide
from pptx.util import Inches, Pt

from megadeck.design_system.tokens import Theme


# ----- Auto-fit / overflow protection -----------------------------------------

# Empirically calibrated against LibreOffice rendering. Inter / SF Pro
# almost never render at width-factor 0.50 — fonts often fall back to a
# wider system family. 0.60 has been verified to over-estimate (which is
# exactly what we want for collision avoidance) on the BMU + Phase-3 demo
# decks across all 36 themes.
_CHAR_WIDTH_FACTOR = 0.60


def fit_title(text: str, *, max_pt: int, min_pt: int = 18, width_in: float = 11.0) -> int:
    """Pick a font size that lets `text` fit in 1-3 lines at width `width_in`.

    Bias-toward-readable: starts at `max_pt`, shrinks by 2pt steps until the
    estimated rendered height stays within 3 lines. We deliberately err on
    the side of *over*-shrinking so titles never overflow into bullets when
    the renderer's font-of-choice isn't installed.
    """
    for size_pt in range(max_pt, min_pt - 1, -2):
        chars_per_line = max(8, int(width_in * 72.0 / (size_pt * _CHAR_WIDTH_FACTOR)))
        n_lines = max(1, math.ceil(len(text) / chars_per_line))
        if n_lines <= 3:
            return size_pt
    return min_pt


def measure_title_height(text: str, *, size_pt: int, width_in: float, line_spacing: float = 1.05) -> float:
    """Estimate rendered height in inches for a title at given pt + width.

    Used to compute where the body should start so it never collides with a
    multi-line title. We err on the side of slightly over-allocating space —
    +0.18in safety pad on top of the worst-case font-fallback line count.
    """
    chars_per_line = max(8, int(width_in * 72.0 / (size_pt * _CHAR_WIDTH_FACTOR)))
    explicit_lines = sum(1 for c in text if c == "\n") + 1
    wrapped_lines = max(1, math.ceil(len(text) / chars_per_line))
    n_lines = max(explicit_lines, wrapped_lines)
    line_height_in = (size_pt * line_spacing) / 72.0
    return n_lines * line_height_in + 0.18  # safety pad for font-fallback drift


# ----- Background --------------------------------------------------------------

def set_slide_bg(slide: Slide, color: RGBColor, theme: Theme | None = None) -> None:
    """Set the slide background.

    For "solid" themes (default) we use the slide's `background` element so
    no extra shape clutters the slide. For gradient themes we add a full-bleed
    rectangle behind everything else and apply a real DrawingML gradient fill.
    The slide's own background stays solid as a fallback for renderers that
    don't support gradients in shape fills (rare).

    Then — regardless of bg style — we layer any `theme.decorations` so orbs
    and ribbons land BEFORE the content. This is what gives a theme like
    Bauhaus or Memphis its visual identity.
    """
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = color

    if theme is not None and theme.bg_style != "solid":
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.util import Emu
        from megadeck.design_system.effects import (
            aurora_background, apply_linear_gradient, apply_radial_gradient,
        )
        rect = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            0, 0,
            Emu(int(theme.slide_width_in * 914400)),
            Emu(int(theme.slide_height_in * 914400)),
        )
        rect.line.fill.background()
        a = theme.bg_aurora_a or theme.bg
        b = theme.bg_aurora_b or theme.bg
        c = theme.bg_aurora_c or theme.bg
        if theme.bg_style == "aurora":
            aurora_background(rect, a=a, b=b, c=c)
        elif theme.bg_style == "vercel-glow":
            apply_radial_gradient(
                rect,
                inner_color=b, outer_color=a,
                inner_alpha=80, outer_alpha=100,
                focus_x=50, focus_y=15,
            )
        elif theme.bg_style == "linear-mesh":
            apply_linear_gradient(
                rect,
                stops=[(0, a, 100), (60, b, 100), (100, c, 100)],
                angle_deg=160.0,
            )

    # Decorations (orbs, mesh, ribbons, geometric, scribble dots…) layer
    # on every theme — including solid-bg ones. This is the difference
    # between a theme and a *design*.
    if theme is not None:
        try:
            from megadeck.design_system.decorations import apply_decorations
            apply_decorations(slide, theme)
        except Exception:
            pass


# ----- Text --------------------------------------------------------------------

def _set_run(
    run,
    *,
    name: Optional[str] = None,
    size_pt: Optional[int] = None,
    bold: Optional[bool] = None,
    italic: Optional[bool] = None,
    color: Optional[RGBColor] = None,
) -> None:
    if name is not None:
        run.font.name = name
    if size_pt is not None:
        run.font.size = Pt(size_pt)
    if bold is not None:
        run.font.bold = bold
    if italic is not None:
        run.font.italic = italic
    if color is not None:
        run.font.color.rgb = color


def add_text(
    slide: Slide,
    *,
    left: float,
    top: float,
    width: float,
    height: float,
    text: str,
    font: str,
    size_pt: int,
    color: RGBColor,
    bold: bool = False,
    italic: bool = False,
    align: Optional[PP_ALIGN] = None,
    anchor: Optional[MSO_ANCHOR] = None,
    line_spacing: Optional[float] = None,
    auto_size: bool = False,
):
    """Add a single-run text box.

    `auto_size` — when True, sets the text frame's auto_size to
    MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT so the box grows vertically to fit its
    content. Use this for titles (where wrap-to-N-lines varies by installed
    fonts) so they never overlap with whatever the renderer placed below.
    """
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    if anchor is not None:
        try:
            tf.vertical_anchor = anchor
        except Exception:
            pass
    if auto_size:
        try:
            from pptx.enum.text import MSO_AUTO_SIZE
            tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        except Exception:
            pass
    p = tf.paragraphs[0]
    if align is not None:
        p.alignment = align
    if line_spacing is not None:
        try:
            p.line_spacing = line_spacing
        except Exception:
            pass
    r = p.add_run()
    r.text = text
    _set_run(r, name=font, size_pt=size_pt, bold=bold, italic=italic, color=color)
    return tb


def add_rich_text(
    slide: Slide,
    *,
    left: float,
    top: float,
    width: float,
    height: float,
    runs: Sequence[Tuple[str, dict]],
    font: str,
    line_spacing: Optional[float] = None,
    align: Optional[PP_ALIGN] = None,
    anchor: Optional[MSO_ANCHOR] = None,
):
    """Compose multiple styled runs in one paragraph."""
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    if anchor is not None:
        try:
            tf.vertical_anchor = anchor
        except Exception:
            pass
    p = tf.paragraphs[0]
    if align is not None:
        p.alignment = align
    if line_spacing is not None:
        try:
            p.line_spacing = line_spacing
        except Exception:
            pass
    for text, opts in runs:
        r = p.add_run()
        r.text = text
        _set_run(
            r,
            name=opts.get("font", font),
            size_pt=opts.get("size_pt"),
            bold=opts.get("bold"),
            italic=opts.get("italic"),
            color=opts.get("color"),
        )
    return tb


# ----- Shapes ------------------------------------------------------------------

def add_rect(
    slide: Slide,
    *,
    left: float,
    top: float,
    width: float,
    height: float,
    fill: Optional[RGBColor] = None,
    line: Optional[RGBColor] = None,
    line_w: Optional[float] = None,
):
    rect = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    if fill is None:
        rect.fill.background()
    else:
        rect.fill.solid()
        rect.fill.fore_color.rgb = fill
    if line is None:
        rect.line.fill.background()
    else:
        rect.line.color.rgb = line
        if line_w is not None:
            rect.line.width = Pt(line_w)
    return rect


def add_round_rect(
    slide: Slide,
    *,
    left: float,
    top: float,
    width: float,
    height: float,
    fill: Optional[RGBColor] = None,
    line: Optional[RGBColor] = None,
    line_w: Optional[float] = None,
    adjust: float = 0.05,
):
    rect = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    rect.adjustments[0] = adjust
    if fill is None:
        rect.fill.background()
    else:
        rect.fill.solid()
        rect.fill.fore_color.rgb = fill
    if line is None:
        rect.line.fill.background()
    else:
        rect.line.color.rgb = line
        if line_w is not None:
            rect.line.width = Pt(line_w)
    return rect


def add_themed_card(
    slide: Slide,
    theme: Theme,
    *,
    left: float,
    top: float,
    width: float,
    height: float,
    accent_border: bool = False,
    adjust: float = 0.05,
):
    """Add a card whose look respects `theme.card_style`.

    flat       → solid surface with hairline + painted-shadow rect behind
    shadow     → solid surface + real DrawingML drop shadow
    glass      → translucent tinted fill + drop shadow (good on aurora bg)
    """
    style = theme.card_style or "flat"
    if style == "flat":
        # Painted shadow underlay
        add_round_rect(
            slide,
            left=left + 0.04, top=top + 0.05,
            width=width, height=height,
            fill=theme.overlay,
            adjust=adjust,
        )
        return add_round_rect(
            slide,
            left=left, top=top, width=width, height=height,
            fill=theme.surface,
            line=theme.accent if accent_border else theme.hairline,
            line_w=1.5 if accent_border else 0.5,
            adjust=adjust,
        )
    if style == "glass":
        rect = add_round_rect(
            slide,
            left=left, top=top, width=width, height=height,
            fill=None, line=None,
            adjust=adjust,
        )
        try:
            from megadeck.design_system.effects import glass_card
            border = theme.accent if accent_border else theme.title
            glass_card(rect, tint=theme.title, border=border, fill_alpha=14)
        except Exception:
            pass
        return rect
    # shadow
    rect = add_round_rect(
        slide,
        left=left, top=top, width=width, height=height,
        fill=theme.surface,
        line=theme.accent if accent_border else theme.hairline,
        line_w=1.5 if accent_border else 0.5,
        adjust=adjust,
    )
    try:
        from megadeck.design_system.effects import apply_drop_shadow
        apply_drop_shadow(
            rect, color="000000", alpha_pct=35, blur_pt=22, distance_pt=4,
        )
    except Exception:
        pass
    return rect


def add_oval(
    slide: Slide,
    *,
    left: float,
    top: float,
    size: float,
    fill: RGBColor,
    line: Optional[RGBColor] = None,
):
    c = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(left), Inches(top), Inches(size), Inches(size),
    )
    c.fill.solid()
    c.fill.fore_color.rgb = fill
    if line is None:
        c.line.fill.background()
    else:
        c.line.color.rgb = line
    return c


def add_v_line(
    slide: Slide,
    *,
    left: float,
    top: float,
    height: float,
    color: RGBColor,
    width: float = 0.04,
):
    return add_rect(
        slide,
        left=left, top=top, width=width, height=height,
        fill=color,
    )


def add_h_line(
    slide: Slide,
    *,
    left: float,
    top: float,
    width: float,
    color: RGBColor,
    height: float = 0.012,
):
    return add_rect(
        slide,
        left=left, top=top, width=width, height=height,
        fill=color,
    )


def add_dot_grid(
    slide: Slide,
    *,
    left: float,
    top: float,
    cols: int = 8,
    rows: int = 3,
    gap: float = 0.13,
    size: float = 0.03,
    color: RGBColor,
) -> None:
    """Draw a tiny dot grid — decorative anchor in a slide corner."""
    for r in range(rows):
        for c in range(cols):
            add_oval(
                slide,
                left=left + c * gap,
                top=top + r * gap,
                size=size,
                fill=color,
            )


# ----- Composite chrome (used by every template) ------------------------------

def add_eyebrow(
    slide: Slide,
    *,
    text: str,
    theme: Theme,
    left: Optional[float] = None,
    top: float = 0.85,
):
    """Standard eyebrow row: thin vertical accent + small uppercase label."""
    x = theme.left_margin_in if left is None else left
    add_v_line(
        slide,
        left=x - 0.10, top=top,
        height=0.30,
        color=theme.accent,
    )
    add_text(
        slide,
        left=x, top=top, width=11.0, height=0.30,
        text=text,
        font=theme.font_body,
        size_pt=theme.type_scale.eyebrow,
        color=theme.muted,
        bold=True,
    )


def add_corner_dotgrid(slide: Slide, theme: Theme) -> None:
    """Tiny decorative dot grid in the top-right corner."""
    add_dot_grid(
        slide,
        left=theme.slide_width_in - 1.95, top=0.35,
        cols=8, rows=3, gap=0.13, size=0.03,
        color=theme.accent_lt,
    )


def add_page_chrome(
    slide: Slide,
    *,
    theme: Theme,
    page_n: int,
    page_total: int,
    section_label: Optional[str] = None,
    on_dark: bool = False,
) -> None:
    """Bottom-of-slide chrome: hairline + section label + page counter."""
    line_color = theme.hairline if not on_dark else theme.overlay
    text_color = theme.light if not on_dark else theme.muted
    add_h_line(
        slide,
        left=theme.left_margin_in, top=7.10,
        width=theme.content_width_in,
        color=line_color,
        height=0.012,
    )
    if section_label:
        add_text(
            slide,
            left=theme.left_margin_in, top=7.18, width=6.0, height=0.30,
            text=section_label,
            font=theme.font_body,
            size_pt=theme.type_scale.eyebrow - 1,
            color=text_color,
            bold=True,
            align=PP_ALIGN.LEFT,
        )
    add_text(
        slide,
        left=theme.slide_width_in - theme.right_margin_in - 1.20,
        top=7.18, width=1.20, height=0.30,
        text=f"{page_n:02d} / {page_total:02d}",
        font=theme.font_body,
        size_pt=theme.type_scale.eyebrow - 1,
        color=text_color,
        align=PP_ALIGN.RIGHT,
    )
