"""Native pptx chart primitives — bar, donut, line.

Uses python-pptx's `slide.shapes.add_chart()` so the output is real
PowerPoint charts (live data, editable in PPT) — not bitmaps.
"""
from __future__ import annotations

from typing import List, Optional, Sequence, Tuple

from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION, XL_LEGEND_POSITION
from pptx.slide import Slide
from pptx.util import Inches, Pt

from megadeck.design_system.tokens import Theme


def add_bar_chart(
    slide: Slide,
    *,
    left_in: float,
    top_in: float,
    width_in: float,
    height_in: float,
    categories: Sequence[str],
    values: Sequence[float],
    series_name: str = "Series",
    theme: Theme,
    horizontal: bool = False,
) -> None:
    """Add a column or bar chart with theme-coloured bars."""
    cd = CategoryChartData()
    cd.categories = list(categories)
    cd.add_series(series_name, list(values))
    chart_type = (
        XL_CHART_TYPE.BAR_CLUSTERED if horizontal else XL_CHART_TYPE.COLUMN_CLUSTERED
    )
    chart = slide.shapes.add_chart(
        chart_type,
        Inches(left_in), Inches(top_in),
        Inches(width_in), Inches(height_in),
        cd,
    ).chart

    chart.has_title = False
    chart.has_legend = False
    plot = chart.plots[0]
    plot.gap_width = 80

    # Style bars with the theme accent
    for series in plot.series:
        fill = series.format.fill
        fill.solid()
        fill.fore_color.rgb = theme.accent
        try:
            series.format.line.fill.background()
        except Exception:
            pass

    # Style category / value axes
    for ax in (chart.category_axis, chart.value_axis):
        try:
            ax.tick_labels.font.size = Pt(10)
            ax.tick_labels.font.color.rgb = theme.body
            ax.format.line.color.rgb = theme.hairline
        except Exception:
            pass


def add_donut_chart(
    slide: Slide,
    *,
    left_in: float,
    top_in: float,
    size_in: float,
    categories: Sequence[str],
    values: Sequence[float],
    theme: Theme,
    palette: Optional[Sequence[RGBColor]] = None,
) -> None:
    """Add a donut chart with one ring of category slices."""
    cd = CategoryChartData()
    cd.categories = list(categories)
    cd.add_series("Donut", list(values))
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.DOUGHNUT,
        Inches(left_in), Inches(top_in),
        Inches(size_in), Inches(size_in),
        cd,
    ).chart

    chart.has_title = False
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.RIGHT
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(10)
    chart.legend.font.color.rgb = theme.body

    # Per-slice colours from a small theme-derived palette.
    if palette is None:
        palette = (
            theme.accent, theme.accent_lt, theme.accent_dk,
            theme.muted, theme.title,
        )
    plot = chart.plots[0]
    points = plot.series[0].points
    for i, pt in enumerate(points):
        f = pt.format.fill
        f.solid()
        f.fore_color.rgb = palette[i % len(palette)]
        try:
            pt.format.line.fill.background()
        except Exception:
            pass


def add_line_chart(
    slide: Slide,
    *,
    left_in: float,
    top_in: float,
    width_in: float,
    height_in: float,
    categories: Sequence[str],
    values: Sequence[float],
    series_name: str = "Series",
    theme: Theme,
) -> None:
    cd = CategoryChartData()
    cd.categories = list(categories)
    cd.add_series(series_name, list(values))
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE,
        Inches(left_in), Inches(top_in),
        Inches(width_in), Inches(height_in),
        cd,
    ).chart
    chart.has_title = False
    chart.has_legend = False
    series = chart.plots[0].series[0]
    line = series.format.line
    line.color.rgb = theme.accent
    line.width = Pt(2.5)


__all__ = ["add_bar_chart", "add_donut_chart", "add_line_chart"]
