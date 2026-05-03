"""Auto-inject a content-aware glass illustration on a slide.

Decorates a slide with:
  * A frosted glass orb in the bottom-right (apple-style highlight)
  * A Lucide icon centered on the orb, picked by content-keyword match
  * A small numerical mark and meta-line for the HUD aesthetic

This runs AFTER the template has rendered the slide, so the illustration
sits on top of content as a deliberate accent — never blocking text.
The orb is positioned in the unused bottom-right region of every layout.
"""
from __future__ import annotations

from typing import Any

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.slide import Slide
from pptx.util import Inches, Pt

from megadeck.design_system.content_icon import icon_for_slide
from megadeck.design_system.effects import (
    apply_radial_gradient,
    apply_solid_fill,
)
from megadeck.design_system.icons import add_icon
from megadeck.design_system.tokens import Theme


def add_glass_illustration(
    slide: Slide,
    slide_data: Any,
    theme: Theme,
    *,
    orb_size_in: float = 1.40,
    margin_in: float = 0.55,
) -> None:
    """Decorate a slide with a content-aware glass orb + Lucide icon.

    The orb is placed in a region of the slide that doesn't conflict with the
    template's content. We pick the position based on (kind, variant):

      * numbered_list/split, numbered_list/timeline → bottom-LEFT
        (the right column has big numbers + items)
      * three_card/asymmetric, kpi_grid/asymmetric  → bottom-LEFT
        (the right side has stacked smaller cards)
      * two_column/vs_arrow                          → skip
        (every region is occupied)
      * everything else                              → bottom-RIGHT
    """
    sw = theme.slide_width_in
    sh = theme.slide_height_in

    kind = getattr(slide_data, "kind", "")
    variant = getattr(slide_data, "variant", None)

    # Slide kinds where the canvas is too dense for an orb anywhere safe:
    if kind == "two_column" and variant == "vs_arrow":
        return
    if kind == "comparison_table":
        return
    # Numbered_list/split with many items: text + numbers fill the whole canvas;
    # any orb position will overlap content. Skip the illustration entirely.
    items_n = len(getattr(slide_data, "items", []) or [])
    if kind == "numbered_list" and variant == "split" and items_n >= 5:
        return
    if kind == "numbered_list" and variant == "timeline" and items_n >= 6:
        return

    # Dense layouts: shrink the orb so it tucks into the corner without
    # colliding with body text. Numbered_list/split fills the entire canvas
    # with text + numbers, so a 1.4in orb is too big.
    items = getattr(slide_data, "items", None)
    n_items = len(items) if items else 0
    is_dense = (kind == "numbered_list" and n_items >= 5) or (
        kind in ("kpi_grid", "bento_grid")
        and len(getattr(slide_data, "tiles", []) or getattr(slide_data, "items", []) or [])
        >= 4
    )
    if is_dense:
        orb_size_in = 0.95

    # Where can it go?
    right_blocked = (
        (kind == "numbered_list" and variant in ("split", "timeline"))
        or (kind == "three_card" and variant == "asymmetric")
        or (kind == "kpi_grid" and variant in ("asymmetric", "data_card"))
        or (kind == "bento_grid" and variant in ("featured", "tiles"))
        or (kind == "photo_card")
    )
    if right_blocked:
        cx = margin_in + orb_size_in / 2 + 0.20
        cy = sh - margin_in - orb_size_in - 0.20
    else:
        cx = sw - margin_in - orb_size_in / 2 - 0.10
        cy = sh - margin_in - orb_size_in - 0.20

    # Outer halo — soft accent radial gradient
    halo_size = orb_size_in * 1.95
    halo = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(cx - halo_size / 2), Inches(cy + orb_size_in / 2 - halo_size / 2),
        Inches(halo_size), Inches(halo_size),
    )
    halo.line.fill.background()
    try:
        apply_radial_gradient(
            halo,
            inner_color=theme.accent, outer_color=theme.accent,
            inner_alpha=24, outer_alpha=0,
        )
    except Exception:
        pass

    # Glass orb — frosted-glass card simulated with low-alpha white fill +
    # thin border + subtle radial gradient.
    orb = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(cx - orb_size_in / 2), Inches(cy),
        Inches(orb_size_in), Inches(orb_size_in),
    )
    try:
        apply_radial_gradient(
            orb,
            inner_color=theme.surface, outer_color=theme.surface,
            inner_alpha=22, outer_alpha=6,
            focus_x=35, focus_y=30,
        )
    except Exception:
        apply_solid_fill(orb, theme.surface, alpha=14)
    try:
        orb.line.color.rgb = theme.accent
        orb.line.width = Pt(0.5)
    except Exception:
        pass

    # Icon centered on the orb
    icon_size = orb_size_in * 0.50
    icon_left = cx - icon_size / 2
    icon_top = cy + (orb_size_in - icon_size) / 2
    icon_name = icon_for_slide(slide_data)
    add_icon(
        slide, icon_name,
        left_in=icon_left, top_in=icon_top,
        size_in=icon_size, color=theme.accent,
    )

    # Tiny meta line under the orb — slide-content tag (e.g. ICON: rocket)
    tag_w = orb_size_in * 1.6
    tag_left = cx - tag_w / 2
    tag_top = cy + orb_size_in + 0.10
    tb = slide.shapes.add_textbox(
        Inches(tag_left), Inches(tag_top), Inches(tag_w), Inches(0.20),
    )
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = f"// {icon_name.replace('-', ' ').upper()}"
    r.font.name = "SF Pro Text"
    r.font.size = Pt(7)
    r.font.color.rgb = theme.muted


__all__ = ["add_glass_illustration"]
