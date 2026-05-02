"""Phase 2: every new slide template renders without errors and produces a
valid pptx file."""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation

from megadeck import render_deck
from megadeck.core.schemas import (
    AgendaItem,
    AgendaSlide,
    BeforeAfterSlide,
    BentoGridSlide,
    BulletItem,
    CardItem,
    CodeSnippetSlide,
    ComparisonTableSlide,
    Deck,
    HeroStatementSlide,
    KpiGridSlide,
    KpiTile,
    NumberedListSlide,
    PullQuoteSlide,
    SectionDividerSlide,
    StepDiagramSlide,
    StepNode,
    TableRow,
    ThreeCardSlide,
    TimelineEvent,
    TimelineSlide,
    TitleSlide,
    TwoColumnSlide,
)


def _build_phase2_sample_deck() -> Deck:
    """Build a deck containing every Phase 1 + Phase 2 slide type."""
    return Deck(
        title="Phase 2 Coverage",
        slides=[
            TitleSlide(
                kind="title",
                eyebrow="Megadeck Phase 2",
                title="Every Template, Every Theme",
                subtitle="Coverage test for the renderer.",
                presenter="Megadeck CI",
                date="2026-05-01",
                venue="GitHub Actions",
            ),
            HeroStatementSlide(
                kind="hero_statement",
                eyebrow="Phase 2",
                statement="Every template renders.",
                supports=["This deck exercises every Phase 1 + 2 slide kind."],
            ),
            AgendaSlide(
                kind="agenda",
                title="Agenda",
                items=[
                    AgendaItem(number="01", title="Foundations", description="Schemas + renderer + theme."),
                    AgendaItem(number="02", title="Phase 2 templates", description="The new ten."),
                    AgendaItem(number="03", title="Animations + import", description="Round-trip + entrances."),
                ],
            ),
            NumberedListSlide(
                kind="numbered_list",
                eyebrow="A long list",
                title="Six items at adaptive sizes",
                items=[BulletItem(head=f"Item {i+1}", tail="Detail") for i in range(6)],
            ),
            ThreeCardSlide(
                kind="three_card",
                eyebrow="Three pillars",
                title="Pillars",
                items=[
                    CardItem(badge=str(i+1).zfill(2), label=f"Pillar {i+1}", description="Description.")
                    for i in range(3)
                ],
            ),
            TwoColumnSlide(
                kind="two_column",
                eyebrow="Two columns",
                title="Stage A vs Stage B",
                left_title="Stage A",
                left_items=[BulletItem(head="A1"), BulletItem(head="A2")],
                right_title="Stage B",
                right_items=[BulletItem(head="B1"), BulletItem(head="B2")],
            ),
            SectionDividerSlide(
                kind="section_divider",
                part_label="Part 2",
                title="Phase 2 templates",
                subtitle="Ten more renderers.",
                dark_background=False,
            ),
            TimelineSlide(
                kind="timeline",
                eyebrow="Roadmap",
                title="Quarterly milestones",
                events=[
                    TimelineEvent(label="Q1", title="Launch", description="MVP shipped."),
                    TimelineEvent(label="Q2", title="Iterate", description="Customer feedback."),
                    TimelineEvent(label="Q3", title="Scale", description="Repeatable channel."),
                    TimelineEvent(label="Q4", title="Expand", description="New markets."),
                ],
            ),
            ComparisonTableSlide(
                kind="comparison_table",
                eyebrow="Plans",
                title="Feature comparison",
                header=["Feature", "Free", "Pro", "Team"],
                rows=[
                    TableRow(cells=["Slides", "10", "Unlimited", "Unlimited"]),
                    TableRow(cells=["Themes", "3", "8", "Custom"]),
                    TableRow(cells=["Critic", "—", "✓", "✓"]),
                ],
            ),
            PullQuoteSlide(
                kind="pull_quote",
                eyebrow="Customer story",
                quote="Megadeck cut my deck-building time from six hours to twenty minutes.",
                author="Anni Zimina",
                role="Co-founder, hAip",
            ),
            BentoGridSlide(
                kind="bento_grid",
                eyebrow="Bento",
                title="Why Megadeck",
                items=[
                    CardItem(badge="01", label=f"Reason {i+1}", description="Card description.")
                    for i in range(4)
                ],
            ),
            KpiGridSlide(
                kind="kpi_grid",
                eyebrow="Metrics",
                title="Q1 results",
                tiles=[
                    KpiTile(label="MRR", value="$24K", delta="+12% MoM", delta_positive=True),
                    KpiTile(label="Users", value="1,832", delta="+220 WoW", delta_positive=True),
                    KpiTile(label="Churn", value="3.1%", delta="-0.4 pts", delta_positive=True),
                    KpiTile(label="NPS", value="52", delta="+4 pts", delta_positive=True),
                ],
            ),
            BeforeAfterSlide(
                kind="before_after",
                eyebrow="The shift",
                title="Before vs After Megadeck",
                before_label="Before",
                before_points=["Manual pptx in Keynote", "Hours per deck", "Inconsistent design"],
                after_label="After",
                after_points=["LLM + DSL", "Minutes per deck", "Audited every slide"],
                verdict="20× faster with the same visual quality.",
            ),
            StepDiagramSlide(
                kind="step_diagram",
                eyebrow="Workflow",
                title="From prompt to pptx",
                steps=[
                    StepNode(title="Prompt", description="Describe the deck."),
                    StepNode(title="Generate", description="LLM emits validated DSL."),
                    StepNode(title="Render", description="DSL → python-pptx."),
                    StepNode(title="Audit", description="Critic flags issues."),
                ],
            ),
            CodeSnippetSlide(
                kind="code_snippet",
                eyebrow="API",
                title="Render a deck in 3 lines",
                language="python",
                code=(
                    "from megadeck import render_deck\n"
                    "deck = generate_deck('My pitch', n_slides=20, theme='default')\n"
                    "render_deck(deck, 'pitch.pptx')"
                ),
                caption="That is the entire surface area of the library.",
            ),
        ],
    )


def test_phase2_renders_for_default_theme(tmp_path: Path) -> None:
    deck = _build_phase2_sample_deck()
    out = tmp_path / "phase2.pptx"
    render_deck(deck, out)
    assert out.exists()
    prs = Presentation(out)
    assert len(prs.slides) == len(deck.slides)


def test_phase2_renders_for_every_theme(tmp_path: Path) -> None:
    base = _build_phase2_sample_deck()
    themes = ["default", "dark", "editorial", "corporate", "linear", "pastel", "neon", "print"]
    for theme in themes:
        deck = base.model_copy(update={"theme": theme})
        out = tmp_path / f"{theme}.pptx"
        render_deck(deck, out)
        assert out.exists(), f"theme={theme} did not produce a pptx"
        # Every theme produces a valid, reopenable file with the same slide count.
        prs = Presentation(out)
        assert len(prs.slides) == len(deck.slides)
