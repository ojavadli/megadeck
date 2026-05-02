"""Phase 3 templates: 11 new templates render and round-trip cleanly across
the 13 registered themes."""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation

from megadeck import render_deck
from megadeck.core.schemas import (
    BulletItem,
    CallToActionSlide,
    Deck,
    FaqItem,
    FaqListSlide,
    FeatureGridSlide,
    FeatureItem,
    LogoGridSlide,
    PricingTableSlide,
    PricingTierData,
    QuestionSlide,
    StatHero,
    StatHeroSlide,
    SwotMatrixSlide,
    TakeawaysSlide,
    TeamGridSlide,
    TeamMember2,
    TestimonialGridSlide,
    TestimonialItem,
    TitleSlide,
)


def _build_phase3_deck() -> Deck:
    return Deck(
        title="Phase 3 Coverage",
        slides=[
            TitleSlide(
                kind="title",
                eyebrow="Megadeck Phase 3",
                title="11 More Templates, 5 More Themes",
                presenter="Megadeck CI",
            ),
            FeatureGridSlide(
                kind="feature_grid",
                eyebrow="Why Megadeck",
                title="Six things you get out of the box",
                features=[
                    FeatureItem(icon_text="◆", title="DSL", description="Constrained schema."),
                    FeatureItem(icon_text="✚", title="Themes", description="13 of them."),
                    FeatureItem(icon_text="↑", title="Critic", description="Visual audit."),
                    FeatureItem(icon_text="◉", title="MCP", description="Claude Code + Cursor."),
                    FeatureItem(icon_text="✦", title="Animations", description="Fade + transitions."),
                    FeatureItem(icon_text="↺", title="Import", description="Round-trip pptx."),
                ],
            ),
            TestimonialGridSlide(
                kind="testimonial_grid",
                eyebrow="Customer Stories",
                title="What people say",
                items=[
                    TestimonialItem(quote="Cut deck-building time 20x.", author="Anni Z.", role="Co-founder"),
                    TestimonialItem(quote="The critic loop is brilliant.", author="Orkhan J.", role="MIT '25"),
                ],
            ),
            TeamGridSlide(
                kind="team_grid",
                eyebrow="Team",
                title="Founding team",
                members=[
                    TeamMember2(initials="OJ", name="Orkhan Javadli", role="Founder"),
                    TeamMember2(initials="AZ", name="Anni Zimina", role="Co-founder"),
                    TeamMember2(initials="MM", name="Megumi Makino", role="Advisor"),
                ],
            ),
            PricingTableSlide(
                kind="pricing_table",
                eyebrow="Pricing",
                title="Simple, scaling",
                tiers=[
                    PricingTierData(
                        name="Free", price="$0", tagline="For trying it out",
                        features=["10 slides / mo", "3 themes", "CLI only"],
                    ),
                    PricingTierData(
                        name="Pro", price="$24 / mo", tagline="For serious users",
                        features=["Unlimited slides", "All themes", "Critic loop", "MCP server"],
                        is_featured=True,
                    ),
                    PricingTierData(
                        name="Team", price="$99 / mo", tagline="For organisations",
                        features=["Everything in Pro", "Custom themes", "Priority support"],
                    ),
                ],
            ),
            TakeawaysSlide(
                kind="takeaways",
                eyebrow="Key takeaways",
                title="Three things to remember",
                items=[
                    BulletItem(head="Constrained DSL beats freeform.", tail="Validate, then render."),
                    BulletItem(head="Themes change look, not content.", tail="One DSL, many decks."),
                    BulletItem(head="Audit before you ship.", tail="Catch overflow before users do."),
                ],
            ),
            CallToActionSlide(
                kind="call_to_action",
                eyebrow="Get started",
                title="Build your next deck in minutes.",
                subtitle="One pip install. One CLI. Or wire it into Claude Code via MCP.",
                button_text="Install Megadeck",
                url="pip install git+https://github.com/ojavadli/megadeck.git",
                footer="MIT-licensed. Open source. Star us.",
            ),
            SwotMatrixSlide(
                kind="swot_matrix",
                eyebrow="Strategy",
                title="Megadeck SWOT",
                strengths=["Open source", "MCP native", "Theme-first"],
                weaknesses=["No web UI yet", "pptx-only"],
                opportunities=["AI-native era", "Cursor + Claude Code adoption"],
                threats=["Closed-source SaaS slide tools"],
            ),
            FaqListSlide(
                kind="faq_list",
                eyebrow="FAQ",
                title="Common questions",
                items=[
                    FaqItem(question="Does it need an LLM key?", answer="Only for `generate`. CLI/import work without one."),
                    FaqItem(question="Can I add custom themes?", answer="Yes. Use `register_theme(Theme(...))`."),
                    FaqItem(question="What about animations?", answer="Native pptx fade / push / morph + entrance staggers."),
                ],
            ),
            StatHeroSlide(
                kind="stat_hero",
                eyebrow="The Result",
                stat=StatHero(value="20×", label="Faster deck-building"),
                title="Hours per deck → minutes per deck.",
                subtitle="Across customer support, scheduling, and pitch decks tested in May 2026.",
            ),
            LogoGridSlide(
                kind="logo_grid",
                eyebrow="Trusted by",
                title="Built and used at",
                logos=["MIT", "Stanford", "Berkeley", "hAip", "OpenAI", "Anthropic", "Google", "BMU"],
            ),
            QuestionSlide(
                kind="question",
                title="Questions?",
                subtitle="Or just star the repo if today helped you.",
                contact="ojavadli@mit.edu",
            ),
        ],
    )


def test_phase3_renders_default_theme(tmp_path: Path) -> None:
    deck = _build_phase3_deck()
    out = tmp_path / "phase3.pptx"
    render_deck(deck, out)
    assert out.exists()
    assert len(Presentation(out).slides) == len(deck.slides)


def test_phase3_renders_every_theme(tmp_path: Path) -> None:
    base = _build_phase3_deck()
    themes = ["default", "dark", "editorial", "corporate", "linear", "pastel",
              "neon", "print", "vibrant", "monochrome", "swiss", "japandi", "glass"]
    for theme in themes:
        deck = base.model_copy(update={"theme": theme})
        out = tmp_path / f"{theme}.pptx"
        render_deck(deck, out)
        prs = Presentation(out)
        assert len(prs.slides) == len(deck.slides)


def test_renderer_dispatch_table_complete() -> None:
    """Every slide kind in the schema is wired into the renderer's dispatch."""
    from megadeck.core.renderer import _RENDERERS
    from megadeck.core import schemas as S
    declared = {
        S.HeroStatementSlide, S.NumberedListSlide, S.ThreeCardSlide, S.TwoColumnSlide,
        S.SectionDividerSlide, S.AgendaSlide, S.TimelineSlide, S.ComparisonTableSlide,
        S.PullQuoteSlide, S.BentoGridSlide, S.KpiGridSlide, S.BeforeAfterSlide,
        S.StepDiagramSlide, S.CodeSnippetSlide, S.TitleSlide, S.FeatureGridSlide,
        S.TestimonialGridSlide, S.TeamGridSlide, S.PricingTableSlide, S.TakeawaysSlide,
        S.CallToActionSlide, S.SwotMatrixSlide, S.FaqListSlide, S.StatHeroSlide,
        S.LogoGridSlide, S.QuestionSlide,
    }
    missing = declared - set(_RENDERERS.keys())
    assert not missing, f"Slide kinds missing from renderer: {missing}"
