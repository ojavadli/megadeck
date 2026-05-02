"""Tests for the DrawingML effects helpers."""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.util import Inches

from megadeck.design_system.effects import (
    apply_drop_shadow,
    apply_glow,
    apply_linear_gradient,
    apply_radial_gradient,
    apply_solid_fill,
    aurora_background,
    glass_card,
    style_card,
)


def _make_blank_slide_with_shape(tmp_path: Path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(1), Inches(1), Inches(4), Inches(2),
    )
    return prs, slide, shape


def _first_sp_xml(pptx_path: Path):
    """Re-open the pptx and return the first <p:sp> element on slide 1."""
    prs = Presentation(pptx_path)
    spTree = prs.slides[0].shapes._spTree
    for child in spTree:
        if child.tag.endswith("}sp"):
            return child
    raise AssertionError("no sp found")


def test_linear_gradient_writes_gradfill(tmp_path: Path) -> None:
    prs, _, sh = _make_blank_slide_with_shape(tmp_path)
    apply_linear_gradient(
        sh,
        stops=[(0, "0E0E20", 100), (50, "1F2D7A", 100), (100, "083344", 100)],
        angle_deg=135.0,
    )
    out = tmp_path / "linear.pptx"
    prs.save(out)
    # Re-open and confirm the gradFill made it in.
    rect = _first_sp_xml(out)
    assert rect.find(qn("p:spPr")).find(qn("a:gradFill")) is not None
    assert rect.find(qn("p:style")) is None  # theme style stripped


def test_radial_gradient_writes_circle_path(tmp_path: Path) -> None:
    prs, _, sh = _make_blank_slide_with_shape(tmp_path)
    apply_radial_gradient(sh, inner_color="3B82F6", outer_color="000000")
    out = tmp_path / "radial.pptx"
    prs.save(out)
    rect = _first_sp_xml(out)
    grad = rect.find(qn("p:spPr")).find(qn("a:gradFill"))
    assert grad is not None
    path = grad.find(qn("a:path"))
    assert path is not None and path.get("path") == "circle"


def test_drop_shadow_adds_outer_shdw(tmp_path: Path) -> None:
    prs, _, sh = _make_blank_slide_with_shape(tmp_path)
    apply_solid_fill(sh, "FFFFFF")
    apply_drop_shadow(sh, blur_pt=10, distance_pt=3, alpha_pct=40)
    out = tmp_path / "shadow.pptx"
    prs.save(out)
    rect = _first_sp_xml(out)
    eff = rect.find(qn("p:spPr")).find(qn("a:effectLst"))
    assert eff is not None
    assert eff.find(qn("a:outerShdw")) is not None


def test_glow_adds_glow_effect(tmp_path: Path) -> None:
    prs, _, sh = _make_blank_slide_with_shape(tmp_path)
    apply_solid_fill(sh, "00E5A0")
    apply_glow(sh, color="00E5A0", radius_pt=20)
    out = tmp_path / "glow.pptx"
    prs.save(out)
    rect = _first_sp_xml(out)
    eff = rect.find(qn("p:spPr")).find(qn("a:effectLst"))
    assert eff is not None
    assert eff.find(qn("a:glow")) is not None


def test_aurora_background_three_stops(tmp_path: Path) -> None:
    prs, _, sh = _make_blank_slide_with_shape(tmp_path)
    aurora_background(sh, a="3B0764", b="BE185D", c="0E7490")
    out = tmp_path / "aurora.pptx"
    prs.save(out)
    rect = _first_sp_xml(out)
    grad = rect.find(qn("p:spPr")).find(qn("a:gradFill"))
    assert grad is not None
    stops = grad.find(qn("a:gsLst")).findall(qn("a:gs"))
    assert len(stops) == 3


def test_glass_card_has_low_alpha_fill(tmp_path: Path) -> None:
    prs, _, sh = _make_blank_slide_with_shape(tmp_path)
    glass_card(sh, tint="FFFFFF", border="FFFFFF", fill_alpha=18)
    out = tmp_path / "glass.pptx"
    prs.save(out)
    rect = _first_sp_xml(out)
    sf = rect.find(qn("p:spPr")).find(qn("a:solidFill"))
    assert sf is not None
    alpha = sf.find(qn("a:srgbClr")).find(qn("a:alpha"))
    assert alpha is not None
    assert int(alpha.get("val")) == 18000  # 18% in OOXML thousandths


def test_style_card_writes_solid_plus_shadow(tmp_path: Path) -> None:
    prs, _, sh = _make_blank_slide_with_shape(tmp_path)
    style_card(sh, fill="FFFFFF", border="E5E7EB", shadow=True)
    out = tmp_path / "style.pptx"
    prs.save(out)
    rect = _first_sp_xml(out)
    sppr = rect.find(qn("p:spPr"))
    assert sppr.find(qn("a:solidFill")) is not None
    assert sppr.find(qn("a:effectLst")).find(qn("a:outerShdw")) is not None


def test_megaslick_themes_register() -> None:
    from megadeck.design_system.tokens import get_theme
    for name in ("motion", "aurora", "vercel", "framer", "linear-pro"):
        t = get_theme(name)
        assert t.bg_style in ("aurora", "vercel-glow", "linear-mesh", "solid")
        assert t.card_style in ("flat", "shadow", "glass")
