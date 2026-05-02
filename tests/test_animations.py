"""Animation primitive tests."""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

from megadeck.animations import apply_transition, stagger_entrance
from megadeck.core.schemas import EntranceKind, TransitionKind


P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"


def _make_slide():
    prs = Presentation()
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    # Add three text shapes so the entrance has something to animate
    for i in range(3):
        tb = slide.shapes.add_textbox(Inches(1), Inches(1 + i), Inches(4), Inches(0.5))
        tb.text_frame.text = f"Block {i + 1}"
    return prs, slide


def test_apply_transition_inserts_xml() -> None:
    _prs, slide = _make_slide()
    apply_transition(slide.element, TransitionKind.FADE)
    assert slide.element.find(f"{{{P_NS}}}transition") is not None


def test_apply_transition_none_removes() -> None:
    _prs, slide = _make_slide()
    apply_transition(slide.element, TransitionKind.FADE)
    apply_transition(slide.element, TransitionKind.NONE)
    assert slide.element.find(f"{{{P_NS}}}transition") is None


def test_stagger_entrance_adds_timing(tmp_path: Path) -> None:
    prs, slide = _make_slide()
    stagger_entrance(slide.element, kind=EntranceKind.FADE_IN, stagger_ms=50)
    timing = slide.element.find(f"{{{P_NS}}}timing")
    assert timing is not None
    # Save to disk and reopen to ensure the XML survives a write/read cycle.
    out = tmp_path / "anim.pptx"
    prs.save(out)
    Presentation(out)


def test_stagger_entrance_none_is_noop() -> None:
    _prs, slide = _make_slide()
    stagger_entrance(slide.element, kind=EntranceKind.NONE)
    assert slide.element.find(f"{{{P_NS}}}timing") is None
