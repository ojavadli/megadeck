"""Tests for `megadeck import`: round-trip an existing pptx into the DSL."""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation

from megadeck import render_deck
from megadeck.core.import_pptx import import_pptx


def test_import_round_trip(tmp_path: Path, sample_deck) -> None:
    src = tmp_path / "src.pptx"
    render_deck(sample_deck, src)
    # Import it back through the heuristic
    imported = import_pptx(src, theme="dark")
    assert imported.theme == "dark"
    assert len(imported.slides) == len(sample_deck.slides)
    # Re-render to confirm the imported deck is renderable.
    out = tmp_path / "round_trip.pptx"
    render_deck(imported, out)
    assert out.exists()
    Presentation(out)
