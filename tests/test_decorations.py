"""Tests for the decoration primitives — orbs, mesh, geometric, scribbles."""
from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches

from megadeck.design_system.decorations import (
    apply_decorations,
    parse_decoration,
    OrbSpec,
    MeshSpec,
    GeometricSpec,
)
from megadeck.design_system.tokens import _THEMES, get_theme


def _blank_slide():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    return prs, slide


def test_parse_orb_spec() -> None:
    spec = parse_decoration({
        "kind": "orb", "x": 0.1, "y": 0.2, "size_in": 5.0,
        "color": "#FF00C8", "alpha_pct": 50,
    })
    assert isinstance(spec, OrbSpec)
    assert spec.size_in == 5.0


def test_parse_mesh_spec() -> None:
    spec = parse_decoration({
        "kind": "mesh",
        "colors": ["#000000", "#222222", "#444444"],
        "angle_deg": 90,
        "alpha_pct": 30,
    })
    assert isinstance(spec, MeshSpec)
    assert len(spec.colors) == 3


def test_parse_geometric_spec() -> None:
    spec = parse_decoration({
        "kind": "geometric", "shape": "triangle",
        "x": 0.5, "y": 0.5, "size_in": 1.0, "color": "#FF0000",
    })
    assert isinstance(spec, GeometricSpec)
    assert spec.shape == "triangle"


def test_parse_unknown_kind_raises() -> None:
    import pytest
    with pytest.raises(ValueError):
        parse_decoration({"kind": "wormhole", "x": 0.5})


def test_every_pool_theme_decorations_render() -> None:
    """Apply every theme's decorations onto a blank slide and assert no
    exceptions, no zero-shape outputs (when decorations were specified)."""
    for name, theme in _THEMES.items():
        if not theme.decorations:
            continue
        prs, slide = _blank_slide()
        before = len(slide.shapes)
        apply_decorations(slide, theme)
        after = len(slide.shapes)
        assert after > before, f"theme {name!r} added 0 decoration shapes"


def test_specific_theme_visual_signatures() -> None:
    """Bauhaus carries 3 geometric primitives; verify they all land."""
    prs, slide = _blank_slide()
    apply_decorations(slide, get_theme("bauhaus"))
    # 3 geometrics → exactly 3 shapes added (no extras)
    assert len(slide.shapes) == 3


def test_decoration_round_trip_via_json(tmp_path: Path) -> None:
    """A theme spec with decorations round-trips JSON ↔ Theme cleanly."""
    from megadeck.design_system.registry import theme_from_dict, theme_to_dict
    spec = {
        "name": "round-trip-decor",
        "bg": "#000000",
        "title": "#FFFFFF",
        "accent": "#FF0000",
        "decorations": [
            {"kind": "orb", "x": 0.5, "y": 0.5, "size_in": 4.0,
             "color": "#FF0000", "alpha_pct": 60},
            {"kind": "edge_ribbon", "edge": "left",
             "color": "#FF0000", "width_in": 0.2},
        ],
    }
    t = theme_from_dict(spec)
    out = theme_to_dict(t)
    assert "decorations" in out
    assert len(out["decorations"]) == 2
    # Round-trip fully through JSON
    txt = json.dumps(out)
    t2 = theme_from_dict(json.loads(txt))
    assert len(t2.decorations) == 2


def test_invalid_decoration_in_pool_json_is_caught(tmp_path: Path) -> None:
    """A theme JSON with a typo'd decoration kind should fail to load."""
    import pytest
    from megadeck.design_system.registry import load_pool_dir
    bad = tmp_path / "pool"
    bad.mkdir()
    (bad / "broken.json").write_text(json.dumps({
        "name": "broken-deco",
        "bg": "#000000",
        "title": "#FFFFFF",
        "accent": "#FF0000",
        "decorations": [{"kind": "wormhole"}],
    }))
    # load_pool_dir doesn't raise — it logs and skips.
    loaded = load_pool_dir(bad)
    assert all(t.name != "broken-deco" for t in loaded)
