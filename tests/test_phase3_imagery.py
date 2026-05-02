"""Tests for Phase A (layouts), B (compositions), C (imagery)."""
from __future__ import annotations

from pathlib import Path

import pytest

from megadeck import render_deck
from megadeck.core.schemas import (
    BarChartSlide,
    BulletItem,
    Deck,
    DonutChartSlide,
    HeroStatementSlide,
    IconBullet,
    IconGridSlide,
    NumberedListSlide,
    PhotoCardSlide,
    TitleSlide,
)
from megadeck.design_system.compositions import (
    COMPOSITIONS,
    DEFAULT_ROTATION,
    apply_composition,
    composition_at_index,
)
from megadeck.design_system.icons import LUCIDE_INLINE, list_icons
from megadeck.design_system.layouts.ingest import ingest_pptx
from megadeck.design_system.layouts.registry import (
    Layout,
    LayoutShape,
    layout_from_dict,
    layout_to_dict,
)


# ---------------------------------------------------------------------------
# Compositions (Phase B)
# ---------------------------------------------------------------------------

def test_all_default_compositions_registered() -> None:
    expected = {
        "typographic", "swiss", "blueprint", "brutalist",
        "editorial", "photographic", "grid", "orbs",
    }
    assert expected.issubset(COMPOSITIONS)


def test_composition_rotation_no_repeat() -> None:
    """The default rotation should never have the same composition back-to-back."""
    seen_at = []
    for i in range(20):
        seen_at.append(composition_at_index(i))
    for prev, nxt in zip(seen_at, seen_at[1:]):
        assert prev != nxt


def test_apply_composition_returns_name() -> None:
    from megadeck.design_system.tokens import get_theme
    from pptx import Presentation
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    name = apply_composition(slide, get_theme("default"), override="swiss")
    assert name == "swiss"


def test_unknown_composition_falls_back_to_typographic() -> None:
    from megadeck.design_system.tokens import get_theme
    from pptx import Presentation
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    name = apply_composition(slide, get_theme("default"), override="not-a-thing")
    assert name == "typographic"


# ---------------------------------------------------------------------------
# Lucide icons (Phase C)
# ---------------------------------------------------------------------------

def test_lucide_inline_30_icons() -> None:
    assert len(LUCIDE_INLINE) >= 30


def test_list_icons_sorted() -> None:
    icons = list_icons()
    assert icons == sorted(icons)
    assert "zap" in icons
    assert "shield" in icons


def test_unknown_icon_falls_back_to_circle() -> None:
    from megadeck.design_system.icons import _icon_svg
    svg = _icon_svg("definitely-not-an-icon")
    assert "<circle" in svg


# ---------------------------------------------------------------------------
# Imagery slide kinds — render end to end (Phase C)
# ---------------------------------------------------------------------------

def test_icon_grid_renders(tmp_path: Path) -> None:
    deck = Deck(
        title="Icon test",
        slides=[
            TitleSlide(kind="title", title="Cover", presenter=""),
            IconGridSlide(
                kind="icon_grid",
                eyebrow="Test",
                title="Icons",
                items=[
                    IconBullet(icon="zap", head="Fast"),
                    IconBullet(icon="shield", head="Safe"),
                    IconBullet(icon="rocket", head="Launch"),
                ],
            ),
        ],
    )
    out = tmp_path / "icons.pptx"
    render_deck(deck, out)
    assert out.exists() and out.stat().st_size > 1000


def test_bar_chart_renders(tmp_path: Path) -> None:
    deck = Deck(
        title="Bar test",
        slides=[
            TitleSlide(kind="title", title="Cover", presenter=""),
            BarChartSlide(
                kind="bar_chart",
                eyebrow="Test",
                title="A bar chart",
                categories=["A", "B", "C"],
                values=[10.0, 25.0, 18.0],
            ),
        ],
    )
    out = tmp_path / "bar.pptx"
    render_deck(deck, out)
    assert out.exists()


def test_donut_chart_renders(tmp_path: Path) -> None:
    deck = Deck(
        title="Donut test",
        slides=[
            TitleSlide(kind="title", title="Cover", presenter=""),
            DonutChartSlide(
                kind="donut_chart",
                eyebrow="Test",
                title="A donut",
                categories=["X", "Y", "Z"],
                values=[40.0, 35.0, 25.0],
            ),
        ],
    )
    out = tmp_path / "donut.pptx"
    render_deck(deck, out)
    assert out.exists()


def test_photo_card_placeholder_renders(tmp_path: Path) -> None:
    """photo=None should render a placeholder card without crashing."""
    deck = Deck(
        title="Photo test",
        slides=[
            TitleSlide(kind="title", title="Cover", presenter=""),
            PhotoCardSlide(
                kind="photo_card",
                eyebrow="Test",
                title="A photo card",
                photo=None,
            ),
        ],
    )
    out = tmp_path / "photo.pptx"
    render_deck(deck, out)
    assert out.exists()


# ---------------------------------------------------------------------------
# Layout ingestion (Phase A)
# ---------------------------------------------------------------------------

def test_layout_dict_roundtrip() -> None:
    lay = Layout(
        name="x", source="test", kind_hint="numbered_list",
        slide_w_in=13.333, slide_h_in=7.5,
        shapes=[
            LayoutShape(role="title", left_in=0.6, top_in=1.2, width_in=10, height_in=1.0),
            LayoutShape(role="number", left_in=0.6, top_in=2.5, width_in=0.6, height_in=0.6),
        ],
    )
    d = layout_to_dict(lay)
    rebuilt = layout_from_dict(d)
    assert rebuilt.name == "x"
    assert len(rebuilt.shapes) == 2
    assert rebuilt.shapes[0].role == "title"


def test_layout_ingestion_produces_layouts(tmp_path: Path) -> None:
    """Render a tiny deck, then ingest it back as layouts."""
    src_deck = Deck(
        title="src",
        slides=[
            TitleSlide(kind="title", title="Cover", presenter="X"),
            NumberedListSlide(
                kind="numbered_list",
                eyebrow="Hi", title="Items",
                items=[BulletItem(head=f"Item {i}", tail="t") for i in range(1, 5)],
            ),
        ],
    )
    src_pptx = tmp_path / "src.pptx"
    render_deck(src_deck, src_pptx)
    layouts = ingest_pptx(src_pptx, save_to=None)
    assert len(layouts) == 2
    names = {l.name for l in layouts}
    assert all(name.startswith("src") for name in names)


# ---------------------------------------------------------------------------
# DSL composition + variant + layout fields
# ---------------------------------------------------------------------------

def test_slide_accepts_composition_and_variant() -> None:
    s = NumberedListSlide(
        kind="numbered_list",
        eyebrow="x", title="t",
        items=[BulletItem(head="a", tail="b"), BulletItem(head="c", tail="d")],
        composition="swiss",
        variant="cards",
    )
    assert s.composition == "swiss"
    assert s.variant == "cards"
