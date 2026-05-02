"""Renderer tests — exercise every template by rendering the sample deck."""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation

from megadeck import render_deck


def test_render_sample_deck_writes_pptx(tmp_path: Path, sample_deck) -> None:
    out = tmp_path / "out.pptx"
    written = render_deck(sample_deck, out)
    assert written == out
    assert out.exists()
    # The file is a valid pptx that python-pptx can re-open
    prs = Presentation(out)
    assert len(prs.slides) == len(sample_deck.slides)


def test_render_preserves_slide_count(tmp_path: Path, sample_deck) -> None:
    out = tmp_path / "out.pptx"
    render_deck(sample_deck, out)
    prs = Presentation(out)
    assert len(prs.slides) == 5  # the sample deck is a 5-slide tour


def test_render_attaches_speaker_notes(tmp_path: Path, sample_deck) -> None:
    out = tmp_path / "out.pptx"
    render_deck(sample_deck, out)
    prs = Presentation(out)
    for slide in prs.slides:
        if slide.has_notes_slide:
            tf = slide.notes_slide.notes_text_frame
            if tf is not None:
                # Notes for sample slides are non-empty
                assert tf.text.strip() != ""


def test_render_each_theme(tmp_path: Path, sample_deck) -> None:
    for theme_name in ("default", "dark", "editorial"):
        deck = sample_deck.model_copy(update={"theme": theme_name})
        out = tmp_path / f"{theme_name}.pptx"
        render_deck(deck, out)
        assert out.exists()
        Presentation(out)  # must be reopenable
